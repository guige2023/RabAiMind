# -*- coding: utf-8 -*-
"""
使用python-pptx直接生成专业PPT - 简化版
"""
import os
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import requests
from io import BytesIO


def create_professional_ppt(
    slides_data: List[Dict],
    output_path: str
) -> str:
    """
    使用python-pptx创建专业PPT
    """
    # 创建16:9的PPT
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    for slide_info in slides_data:
        slide_type = slide_info.get("slide_type", "content")
        title = slide_info.get("title", "")
        content = slide_info.get("content", [])
        image_url = slide_info.get("image_url", "")
        
        # 创建空白幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 根据类型创建不同布局
        if slide_type == "title" or slide_type == "title_slide":
            _create_title_slide(prs, slide, title, content, image_url)
        elif slide_type == "thank_you":
            _create_thank_you_slide(slide, title, content)
        else:
            # 内容页
            _create_content_slide(prs, slide, title, content, image_url)
    
    # 保存
    prs.save(output_path)
    return output_path


def _create_title_slide(prs, slide, title: str, content: List, image_url: str):
    """创建封面页 - 全屏图片+文字居中"""
    # 全屏图片
    img_path = _load_image(image_url)
    if img_path:
        try:
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            with Image.open(img_path) as img:
                img_width, img_height = img.size
                ratio = max(slide_width / img_width, slide_height / img_height)
                new_width = img_width * ratio
                new_height = img_height * ratio
            
            left = (slide_width - new_width) / 2
            top = (slide_height - new_height) / 2
            
            slide.shapes.add_picture(img_path, left, top, width=new_width, height=new_height)
        except Exception as e:
            print(f"添加图片失败: {e}")
    
    # 半透明遮罩
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shape.fill.transparency = 0.4
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if content:
        sub_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(12), Inches(1))
        tf = sub_box.text_frame
        for i, item in enumerate(content[:3]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(220, 220, 220)
            p.alignment = PP_ALIGN.CENTER


def _create_content_slide(prs, slide, title: str, content: List, image_url: str):
    """创建内容页 - 全屏图片+底部文字"""
    # 全屏图片
    img_path = _load_image(image_url)
    if img_path:
        try:
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            with Image.open(img_path) as img:
                img_width, img_height = img.size
                ratio = max(slide_width / img_width, slide_height / img_height)
                new_width = img_width * ratio
                new_height = img_height * ratio
            
            left = (slide_width - new_width) / 2
            top = (slide_height - new_height) / 2
            
            slide.shapes.add_picture(img_path, left, top, width=new_width, height=new_height)
        except Exception as e:
            print(f"添加图片失败: {e}")
    
    # 底部半透明遮罩
    shape = slide.shapes.add_shape(1, 0, Inches(6.5), prs.slide_width, Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shape.fill.transparency = 0.6
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(15), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 内容要点 - 横向排列
    if content:
        for i, item in enumerate(content[:3]):
            x = Inches(0.5) + i * Inches(5.2)
            content_box = slide.shapes.add_textbox(x, Inches(7.8), Inches(5), Inches(1))
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "▸ " + item
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(200, 200, 200)


def _create_thank_you_slide(prs, slide, title: str, content: List):
    """创建尾页"""
    # 深蓝色背景
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(22, 93, 255)
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 联系方式
    if content:
        sub_box = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(10), Inches(1))
        tf = sub_box.text_frame
        for i, item in enumerate(content[:2]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.alignment = PP_ALIGN.CENTER


def _load_image(url: str):
    """从URL加载图片并保存为临时文件"""
    try:
        response = requests.get(url, timeout=30)
        if response.status_code == 200:
            # 保存为临时文件
            import tempfile
            import os
            
            # 确定文件扩展名
            content_type = response.headers.get('Content-Type', '')
            ext = '.jpg'
            if 'png' in content_type:
                ext = '.png'
            elif 'gif' in content_type:
                ext = '.gif'
            
            # 创建临时文件
            fd, path = tempfile.mkstemp(suffix=ext)
            os.write(fd, response.content)
            os.close(fd)
            
            return path
    except Exception as e:
        print(f"下载图片失败: {url}, {e}")
    return None


# 全局变量
prs = None
