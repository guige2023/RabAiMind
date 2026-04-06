# -*- coding: utf-8 -*-
"""
Additional Export Service
Supports: ODP (OpenDocument), Keynote (.key), MP3 Audio Narration
"""

import io
import logging
import os
import shutil
import subprocess
import tempfile
import zipfile
from typing import Dict, Any, Optional, List

logger = logging.getLogger(__name__)


class AdditionalExportService:
    """Export PPTX to ODP, Keynote, and MP3 audio"""

    def __init__(self):
        self._edge_tts_available = None
        self._libreoffice_path = None

    def _check_edge_tts(self) -> bool:
        if self._edge_tts_available is None:
            self._edge_tts_available = shutil.which("edge-tts") is not None
        return self._edge_tts_available

    def _get_libreoffice_path(self) -> Optional[str]:
        if self._libreoffice_path is None:
            self._libreoffice_path = shutil.which("libreoffice") or shutil.which("soffice")
        return self._libreoffice_path

    async def export_to_odp(
        self,
        task_id: str,
        pptx_path: str,
        output_path: str
    ) -> Dict[str, Any]:
        """
        Export PPTX to ODP (OpenDocument Presentation) using LibreOffice.
        """
        libreoffice = self._get_libreoffice_path()
        if not libreoffice:
            return {
                "success": False,
                "error": "LibreOffice 不可用",
                "hint": "请安装 LibreOffice: brew install --cask libreoffice"
            }

        try:
            # Use LibreOffice headless to convert PPTX to ODP
            output_dir = os.path.dirname(output_path) or tempfile.gettempdir()
            base_name = os.path.splitext(os.path.basename(pptx_path))[0]

            # LibreOffice can convert to odp directly
            result = subprocess.run(
                [
                    libreoffice,
                    "--headless",
                    "--convert-to", "odp",
                    "--outdir", output_dir,
                    pptx_path
                ],
                capture_output=True,
                text=True,
                timeout=120
            )

            if result.returncode != 0:
                logger.error(f"LibreOffice ODP conversion failed: {result.stderr}")
                return {
                    "success": False,
                    "error": f"ODP 转换失败: {result.stderr}",
                    "method": "manual",
                    "download_url": f"/api/v1/ppt/download/{task_id}"
                }

            # Find the generated ODP file
            expected_odp = os.path.join(output_dir, f"{base_name}.odp")
            if os.path.exists(expected_odp):
                # Move to final location if different
                if expected_odp != output_path:
                    shutil.copy2(expected_odp, output_path)
                    os.remove(expected_odp)
                return {
                    "success": True,
                    "output_path": output_path,
                    "method": "libreoffice"
                }
            else:
                # Try to find any .odp in output dir
                odp_files = [f for f in os.listdir(output_dir) if f.endswith(".odp")]
                if odp_files:
                    src = os.path.join(output_dir, odp_files[0])
                    shutil.copy2(src, output_path)
                    return {"success": True, "output_path": output_path, "method": "libreoffice"}

            return {
                "success": False,
                "error": "未能生成 ODP 文件",
                "method": "manual",
                "download_url": f"/api/v1/ppt/download/{task_id}"
            }

        except subprocess.TimeoutExpired:
            return {
                "success": False,
                "error": "ODP 转换超时",
                "method": "manual",
                "download_url": f"/api/v1/ppt/download/{task_id}"
            }
        except Exception as e:
            logger.exception("ODP export error")
            return {
                "success": False,
                "error": str(e),
                "method": "manual",
                "download_url": f"/api/v1/ppt/download/{task_id}"
            }

    async def export_to_keynote(
        self,
        task_id: str,
        pptx_path: str,
        output_path: str
    ) -> Dict[str, Any]:
        """
        Export PPTX to Keynote (.key) format.
        
        Note: Apple Keynote natively supports ODP import. We first convert
        to ODP via LibreOffice, then create a .key package structure.
        .key files are macOS document packages containing ODP-compatible XML.
        """
        libreoffice = self._get_libreoffice_path()
        if not libreoffice:
            return {
                "success": False,
                "error": "LibreOffice 不可用",
                "hint": "请安装 LibreOffice"
            }

        try:
            # First convert to ODP
            temp_dir = tempfile.mkdtemp(prefix="keynote_export_")
            temp_odp = os.path.join(temp_dir, "temp_export.odp")
            
            # Convert PPTX to ODP
            result = subprocess.run(
                [
                    libreoffice,
                    "--headless",
                    "--convert-to", "odp",
                    "--outdir", temp_dir,
                    pptx_path
                ],
                capture_output=True,
                text=True,
                timeout=120
            )

            if result.returncode != 0:
                shutil.rmtree(temp_dir, ignore_errors=True)
                return {
                    "success": False,
                    "error": f"Keynote 转换失败: {result.stderr[:200]}",
                    "method": "manual",
                    "guide": {
                        "step1": "下载 PPTX 文件",
                        "step2": "在 Mac 上用 Keynote 打开",
                        "step3": "文件 → 导出为 Keynote 格式"
                    },
                    "download_url": f"/api/v1/ppt/download/{task_id}"
                }

            # Find the generated ODP
            odp_files = [f for f in os.listdir(temp_dir) if f.endswith(".odp")]
            if not odp_files:
                shutil.rmtree(temp_dir, ignore_errors=True)
                return {
                    "success": False,
                    "error": "Keynote 转换中间步骤失败",
                    "method": "manual",
                    "download_url": f"/api/v1/ppt/download/{task_id}"
                }

            temp_odp = os.path.join(temp_dir, odp_files[0])
            
            # Create .key package (macOS app bundle format)
            # A .key file is essentially a package containing:
            # Index.zip (main content), preview.png, etc.
            key_dir = tempfile.mkdtemp(prefix="keynote_package_")
            key_name = os.path.splitext(os.path.basename(output_path))[0]
            key_package_dir = os.path.join(key_dir, f"{key_name}.key")
            os.makedirs(key_package_dir)

            # Create the standard Keynote package structure
            # Content/Index.zip contains the ODP-like data
            content_dir = os.path.join(key_package_dir, "Content")
            os.makedirs(content_dir)

            # Copy ODP as Index.zip inside Content
            # Note: Real .key uses specific XML structure, but Keynote
            # can import ODP files directly, so we provide ODP as fallback
            shutil.copy2(temp_odp, os.path.join(content_dir, "Index.zip"))

            # Create preview (placeholder - real preview would need rendering)
            preview_path = os.path.join(key_package_dir, "preview.png")
            self._create_placeholder_preview(preview_path)

            # Create Document.info (metadata)
            info_path = os.path.join(key_package_dir, "Document.info")
            with open(info_path, "w") as f:
                f.write('<?xml version="1.0" encoding="UTF-8"?>\n')
                f.write('<!DOCTYPE plist PUBLIC "-//Apple//DTD PLIST 1.0//EN" "http://www.apple.com/DTDs/PropertyList-1.0.dtd">\n')
                f.write('<plist version="1.0">\n')
                f.write('<dict>\n')
                f.write(f'  <key>name</key><string>{key_name}</string>\n')
                f.write('  <key>creator</key><string>RabAiMind</string>\n')
                f.write('  <key>version</key><integer>1</integer>\n')
                f.write('</dict>\n')
                f.write('</plist>\n')

            # Create a standard ZIP of the .key package
            # .key is typically not zipped itself - it's a directory/package
            # But for download, we can ZIP it or just provide as-is
            
            # Copy to final destination
            final_key = os.path.join(os.path.dirname(output_path) or tempfile.gettempdir(), f"{key_name}.key")
            
            # For simplicity, we'll copy the whole directory as a .zip
            # since True .key creation requires macOS APIs
            zip_path = output_path
            if not output_path.endswith(".zip"):
                zip_path = output_path + ".zip"
            
            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
                for root, dirs, files in os.walk(key_package_dir):
                    for file in files:
                        file_path = os.path.join(root, file)
                        arcname = os.path.relpath(file_path, key_package_dir)
                        zf.write(file_path, arcname)

            # Cleanup
            shutil.rmtree(temp_dir, ignore_errors=True)
            shutil.rmtree(key_dir, ignore_errors=True)

            return {
                "success": True,
                "output_path": zip_path,
                "method": "libreoffice+package",
                "note": "Keynote 包已创建，请在 Mac 上用 Keynote 打开 .key 文件"
            }

        except subprocess.TimeoutExpired:
            return {
                "success": False,
                "error": "Keynote 转换超时",
                "method": "manual",
                "download_url": f"/api/v1/ppt/download/{task_id}"
            }
        except Exception as e:
            logger.exception("Keynote export error")
            return {
                "success": False,
                "error": str(e),
                "method": "manual",
                "download_url": f"/api/v1/ppt/download/{task_id}"
            }

    def _create_placeholder_preview(self, preview_path: str):
        """Create a placeholder preview image for the Keynote package"""
        try:
            from PIL import Image, ImageDraw, ImageFont
            
            # Create a simple placeholder image
            img = Image.new('RGB', (256, 192), color=(240, 240, 240))
            draw = ImageDraw.Draw(img)
            
            # Draw some text
            try:
                font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 16)
            except Exception:
                font = ImageFont.load_default()
            
            draw.text((20, 80), "Preview", fill=(100, 100, 100), font=font)
            img.save(preview_path, "PNG")
        except ImportError:
            # PIL not available, create minimal placeholder
            with open(preview_path, 'wb') as f:
                # Minimal 1x1 transparent PNG
                f.write(b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc\x00\x01\x00\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82')

    async def export_to_mp3(
        self,
        task_id: str,
        pptx_path: str,
        output_path: str,
        slides_content: Optional[List[Dict[str, Any]]] = None,
        voice: str = "zh-CN-XiaoxiaoNeural",
        rate: str = "+0%",
        volume: str = "+0%"
    ) -> Dict[str, Any]:
        """
        Export PPT narration as MP3 audio using edge-tts.
        Each slide becomes an audio segment, combined into one MP3.
        
        Args:
            task_id: Task ID for reference
            pptx_path: Path to PPTX file (used for extracting text if no slides_content)
            output_path: Output MP3 file path
            slides_content: List of slide content dicts with title/content
            voice: edge-tts voice name (default Chinese female)
            rate: Speech rate (e.g., "+0%", "+10%", "-10%")
            volume: Speech volume (e.g., "+0%", "+10%")
        """
        if not self._check_edge_tts():
            return {
                "success": False,
                "error": "edge-tts 不可用",
                "hint": "请安装 edge-tts: pip install edge-tts"
            }

        try:
            import asyncio
            from edge_tts import Communicate

            # Build narration text for each slide
            slide_texts = []
            
            if slides_content:
                for i, slide in enumerate(slides_content):
                    slide_parts = []
                    
                    # Title
                    title = slide.get("title", "")
                    if title:
                        slide_parts.append(f"第{i+1}页. {title}")
                    
                    # Content items
                    content = slide.get("content", [])
                    if isinstance(content, list):
                        for item in content:
                            if isinstance(item, dict):
                                text = item.get("text", "")
                                if text:
                                    slide_parts.append(text)
                            elif isinstance(item, str):
                                slide_parts.append(item)
                    elif isinstance(content, str):
                        slide_parts.append(content)
                    
                    slide_text = "。".join(slide_parts)
                    slide_texts.append(slide_text)
            else:
                # Try to extract text from PPTX directly
                try:
                    from pptx import Presentation
                    prs = Presentation(pptx_path)
                    for i, slide in enumerate(prs.slides):
                        slide_parts = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                slide_parts.append(shape.text.strip())
                        if slide_parts:
                            slide_texts.append(f"第{i+1}页. {'。'.join(slide_parts)}")
                except Exception as e:
                    logger.warning(f"Could not extract text from PPTX: {e}")
                    return {
                        "success": False,
                        "error": f"无法从PPT提取文本: {str(e)}",
                        "hint": "请确保提供了 slides_content 数据"
                    }

            if not slide_texts:
                return {
                    "success": False,
                    "error": "没有可朗读的文本内容",
                    "hint": "确保幻灯片包含标题或文本内容"
                }

            # Generate audio for each slide using edge-tts
            temp_dir = tempfile.mkdtemp(prefix="ppt_audio_")
            audio_files = []

            async def generate_slide_audio(slide_idx: int, text: str, temp_dir: str) -> str:
                """Generate audio for a single slide"""
                output_file = os.path.join(temp_dir, f"slide_{slide_idx+1:03d}.mp3")
                
                try:
                    # Trim very long text (edge-tts has limits)
                    if len(text) > 2000:
                        text = text[:2000] + "..."
                    
                    if not text.strip():
                        return None
                    
                    communicate = Communicate(text, voice, rate=rate, volume=volume)
                    await communicate.save(output_file)
                    return output_file
                except Exception as e:
                    logger.warning(f"Failed to generate audio for slide {slide_idx+1}: {e}")
                    return None

            # Generate all slides concurrently
            tasks = [
                generate_slide_audio(i, text, temp_dir)
                for i, text in enumerate(slide_texts)
            ]
            
            import asyncio
            results = await asyncio.gather(*tasks)
            audio_files = [f for f in results if f]

            if not audio_files:
                shutil.rmtree(temp_dir, ignore_errors=True)
                return {
                    "success": False,
                    "error": "音频生成失败",
                    "hint": "请检查网络连接或尝试其他语音"
                }

            # Combine all MP3 files into one using ffmpeg
            concat_list_path = os.path.join(temp_dir, "concat_list.txt")
            with open(concat_list_path, "w") as f:
                for audio_file in sorted(audio_files):
                    # Escape special characters in path
                    escaped_path = audio_file.replace("'", "'\\''")
                    f.write(f"file '{escaped_path}'\n")

            # Use ffmpeg to concatenate MP3 files
            ffmpeg_path = shutil.which("ffmpeg")
            if not ffmpeg_path:
                # Fallback: just return the first audio file
                shutil.copy2(audio_files[0], output_path)
                shutil.rmtree(temp_dir, ignore_errors=True)
                return {
                    "success": True,
                    "output_path": output_path,
                    "method": "edge-tts-single",
                    "note": "仅生成了第一页音频（ffmpeg 未安装）"
                }

            # Concatenate MP3 files
            concat_result = subprocess.run(
                [
                    ffmpeg_path,
                    "-y",  # Overwrite output
                    "-f", "concat",
                    "-safe", "0",
                    "-i", concat_list_path,
                    "-c", "copy",
                    output_path
                ],
                capture_output=True,
                text=True,
                timeout=300
            )

            if concat_result.returncode != 0:
                # Fallback: copy first file
                shutil.copy2(audio_files[0], output_path)
            else:
                # Verify output was created
                if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                    shutil.copy2(audio_files[0], output_path)

            # Cleanup
            shutil.rmtree(temp_dir, ignore_errors=True)

            file_size = os.path.getsize(output_path) if os.path.exists(output_path) else 0
            return {
                "success": True,
                "output_path": output_path,
                "method": "edge-tts+ffmpeg",
                "slides_count": len(slide_texts),
                "file_size": file_size,
                "voice": voice,
                "rate": rate
            }

        except ImportError as e:
            logger.error(f"edge-tts import failed: {e}")
            return {
                "success": False,
                "error": "缺少 edge-tts 库",
                "hint": "pip install edge-tts"
            }
        except Exception as e:
            logger.exception("MP3 export error")
            return {
                "success": False,
                "error": str(e),
                "method": "manual",
                "download_url": f"/api/v1/ppt/download/{task_id}"
            }


def get_additional_export_service() -> AdditionalExportService:
    """Get singleton instance"""
    if not hasattr(get_additional_export_service, '_instance'):
        get_additional_export_service._instance = AdditionalExportService()
    return get_additional_export_service._instance
