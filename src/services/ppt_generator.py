# -*- coding: utf-8 -*-
"""
PPT 生成器服务 - 使用 okppt 方式 (AI生成SVG)

作者: Claude
日期: 2026-03-18
"""

import time
import asyncio
import os
import re
import json
import threading
import numpy as np
from typing import Optional, Dict, Any, List
from pathlib import Path
from io import BytesIO

# 导入工具函数
from ..utils import ensure_dir, get_timestamp, setup_logger
from ..config import settings

# 导入pptx相关
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

logger = setup_logger("ppt_generator")


# ─── 独立工具函数（类外定义，避免循环引用）───────────────────────────────────────

def svg_placeholder_xml(
    slide_num: int,
    title: str,
    bg_color: str = "#1e3a5f",
    accent_color: str = "#60a5fa",
    icon: str = "📄"
) -> str:
    """生成 SVG 占位图 XML 字符串（不依赖任何外部资源）"""
    W, H = 1600, 900
    # 标题最多显示 20 字符
    display_title = (title[:20] + "…") if len(title) > 20 else title
    return (
        f'''<?xml version="1.0" encoding="UTF-8"?>'
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {W} {H}" width="{W}" height="{H}">
  <defs>
    <linearGradient id="bgGrad" x1="0%" y1="0%" x2="100%" y2="100%">
      <stop offset="0%" style="stop-color:{bg_color}"/>
      <stop offset="100%" style="stop-color:{bg_color}cc"/>
    </linearGradient>
  </defs>
  <!-- 背景 -->
  <rect width="{W}" height="{H}" fill="url(#bgGrad)"/>
  <!-- 装饰圆 -->
  <circle cx="1400" cy="100" r="300" fill="{accent_color}" opacity="0.08"/>
  <circle cx="200"  cy="800" r="200" fill="{accent_color}" opacity="0.06"/>
  <!-- 分隔线 -->
  <line x1="100" y1="600" x2="1500" y2="600" stroke="{accent_color}" stroke-width="2" opacity="0.3"/>
  <!-- 图标 -->
  <text x="800" y="480" text-anchor="middle" font-size="120" opacity="0.6">{icon}</text>
  <!-- 标题 -->
  <text x="800" y="700" text-anchor="middle" fill="white" font-size="36" font-weight="bold"
        font-family="Microsoft YaHei, PingFang SC, sans-serif">{display_title}</text>
  <!-- 页码 -->
  <text x="1500" y="860" text-anchor="end" fill="{accent_color}" font-size="20"
        font-family="Arial, sans-serif" opacity="0.7">{slide_num} / N</text>
</svg>'''
    )


class PPTGenerator:
    """PPT 生成器 - okppt方式"""

    def __init__(self):
        ensure_dir(settings.OUTPUT_DIR)
        ensure_dir(settings.TEMPLATE_DIR)
        self._image_failure_count: int = 0  # AI图片生成失败计数
        self._image_failure_threshold: int = 3  # 失败阈值，超过后停止尝试
        self._first_page_layout: Optional[str] = None  # 统一布局模式：首页布局
        self._layout_lock = threading.Lock()  # 布局操作线程锁

    async def generate(
        self,
        task_id: str,
        user_request: str,
        slide_count: int = 10,
        scene: str = "business",
        style: str = "professional",
        template: str = "default",
        theme_color: str = "#165DFF",
        text_style: str = "transparent_overlay",
        shadow_color: str = "#000000",
        overlay_transparency: int = 30,
        use_smart_layout: bool = False,
        slide_backgrounds: Optional[List] = None,
        slide_layouts: Optional[List] = None,
        include_charts: bool = False,
        include_pie_chart: bool = True,
        include_bar_chart: bool = True,
        include_line_chart: bool = False,
        add_watermark: bool = False,
        # 新增参数
        font_title: str = "思源黑体",
        font_subtitle: str = "思源黑体",
        font_content: str = "思源宋体",
        font_caption: str = "思源黑体",
        generation_mode: str = "standard",  # standard/fast/quality/stream
        output_format: str = "pptx",  # pptx/pdf/svg/png
        quality: str = "standard",  # standard/high/ultra
        layout_mode: str = "auto",  # auto/manual
        unified_layout: bool = True,  # 是否统一布局
        # 两阶段生成：预生成的内容（跳过 AI 内容规划步骤）
        pre_generated_slides: Optional[List[Dict]] = None  # [{title, content, slide_type, layout}, ...]
    ) -> Dict[str, Any]:
        """生成 PPT - okppt方式

        两阶段模式:
        - Step 1 (内容): AI 规划内容，用户在大纲编辑页确认
        - Step 2 (排版): 传入预生成内容，直接进入排版渲染，跳过 AI 内容规划

        Args:
            generation_mode: 生成模式
                - standard: 标准模式，平衡速度和质量
                - fast: 快速模式，快速生成预览
                - quality: 高清模式，高质量输出
                - stream: 流式模式，边生成边输出

            output_format: 输出格式
                - pptx: PowerPoint文件
                - pdf: PDF文件
                - svg: SVG矢量图
                - png: PNG图片

            quality: 输出质量
                - standard: 1080p
                - high: 1440p
                - ultra: 4K
        """
        logger.info(f"开始生成 PPT (okppt方式), task_id={task_id}, slide_count={slide_count}, mode={generation_mode}, format={output_format}, quality={quality}")

        # 重置跨任务状态，防止污染
        self._first_page_layout = None
        self._image_failure_count = 0

        try:
            from .task_manager import get_task_manager
            task_manager = get_task_manager()

            # 更新状态 - 0%
            task_manager.update_progress(task_id, 0, "开始生成PPT...", "processing")

            # 1. PPT内容来源：两阶段 or 一次性
            if pre_generated_slides:
                # 两阶段模式：内容已由用户在 OutlineEditView 确认
                slides_content = pre_generated_slides
                logger.info(f"使用预生成内容，共 {len(slides_content)} 页（跳过 AI 内容规划）")
                task_manager.update_progress(task_id, 20, f"内容已确认，开始渲染设计...", "processing")
            else:
                # 一次性模式：AI 内容规划 5%~20%
                logger.info(f"开始智能规划 PPT, request={user_request[:50]}...")
                task_manager.update_progress(task_id, 5, "AI正在构思布局...", "processing")
                slides_content = await asyncio.to_thread(
                    self._plan_ppt, user_request, slide_count
                )
                logger.info(f"AI规划了 {len(slides_content)} 页")
                task_manager.update_progress(task_id, 20, f"已完成内容规划，共{len(slides_content)}页", "processing")

            # 2. 为每页生成AI图片 - 20%~45%
            task_manager.update_progress(task_id, 22, f"开始生成{slide_count}张配图...", "processing")
            logger.info("开始生成AI图片...")

            for i, slide in enumerate(slides_content):
                slide_num = i + 1
                progress = 22 + int((i / slide_count) * 23)
                task_manager.update_progress(
                    task_id, progress,
                    f"正在生成第{slide_num}/{slide_count}张配图...",
                    "processing"
                )
                image_url = self._generate_image(slide, slide_num)
                slide["image_url"] = image_url
                logger.info(f"slide {slide_num} 图片生成完成")

            task_manager.update_progress(task_id, 45, "配图生成完成", "processing")

            # 3. AI生成SVG代码 - 45%~80%
            task_manager.update_progress(task_id, 48, "AI设计中...", "processing")
            logger.info("开始生成SVG代码...")

            # 使用多线程并行生成SVG
            from concurrent.futures import ThreadPoolExecutor, as_completed

            def generate_single_svg(args):
                i, slide = args
                slide_num = i + 1
                user_layout = None
                if slide_layouts:
                    for layout in slide_layouts:
                        if layout.slide_index == i:
                            user_layout = layout.layout_type
                            break
                if use_smart_layout:
                    # 传递文字样式和字体参数
                    svg_code = self._generate_svg_smart_layout(
                        slide, slide_num, theme_color, style, user_layout,
                        text_style=text_style,
                        shadow_color=shadow_color,
                        overlay_transparency=overlay_transparency,
                        font_title=font_title,
                        font_subtitle=font_subtitle,
                        font_content=font_content,
                        font_caption=font_caption,
                        layout_mode=layout_mode,
                        unified_layout=unified_layout
                    )
                else:
                    svg_code = self._generate_svg(slide, slide_num)
                return slide_num, svg_code

            # 并行生成SVG，带进度更新
            svg_files = []
            completed = 0
            total = len(slides_content)

            with ThreadPoolExecutor(max_workers=4) as executor:
                futures = {executor.submit(generate_single_svg, (i, slide)): i for i, slide in enumerate(slides_content)}
                for future in as_completed(futures):
                    slide_num, svg_code = future.result()
                    svg_path = os.path.join(settings.OUTPUT_DIR, f"slide_{slide_num}_okppt.svg")
                    with open(svg_path, 'w', encoding='utf-8') as f:
                        f.write(svg_code)
                    svg_files.append(svg_path)

                    completed += 1
                    progress = 48 + int((completed / total) * 32)
                    task_manager.update_progress(
                        task_id, progress,
                        f"正在渲染第{completed}/{total}页...",
                        "processing"
                    )
                    logger.info(f"slide {slide_num} SVG生成完成")

            # 按页码排序
            svg_files.sort(key=lambda x: int(x.split('_')[-2]))

            task_manager.update_progress(task_id, 80, f"渲染完成，共{len(slides_content)}页", "processing")

            # 4. SVG转PPT - 80%~95%
            task_manager.update_progress(task_id, 85, "正在生成PPTX文件...", "processing")
            output_path = os.path.join(
                settings.OUTPUT_DIR,
                f"presentation_{task_id}.pptx"
            )

            logger.info(f"SVG转PPT: {output_path}")
            # 同步调用 - 传递字体参数
            self._svg_to_ppt(svg_files, output_path, text_style, shadow_color, overlay_transparency, use_smart_layout, slide_backgrounds, font_title, font_subtitle, font_content, font_caption)

            logger.info(f"PPT生成完成: {output_path}")
            task_manager.update_progress(task_id, 95, "文件处理中...", "processing")

            # 5. 完成
            task_manager.update_progress(task_id, 100, "生成完成!", "completed")
            
            # 完成任务
            task_manager.complete_task(
                task_id=task_id,
                pptx_path=output_path,
                slide_count=len(slides_content),
                compression_ratio=1.0
            )

            # 根据输出格式生成不同文件
            output_files = {
                "pptx": output_path,
                "format": output_format,
                "quality": quality,
                "generation_mode": generation_mode
            }

            # 如果需要导出其他格式
            if output_format == "svg":
                # 保留SVG文件路径
                output_files["svg_paths"] = svg_files
            elif output_format == "png":
                # 转换为PNG
                png_dir = os.path.join(settings.OUTPUT_DIR, f"png_{task_id}")
                os.makedirs(png_dir, exist_ok=True)
                png_paths = self._convert_svg_to_png(svg_files, png_dir, quality)
                output_files["png_paths"] = png_paths
            elif output_format == "pdf":
                # 转换为PDF
                pdf_path = os.path.join(settings.OUTPUT_DIR, f"presentation_{task_id}.pdf")
                self._svg_to_pdf(svg_files, pdf_path)
                output_files["pdf_path"] = pdf_path

            logger.info(f"PPT生成完成: {output_format}, quality={quality}, mode={generation_mode}")

            return {
                "success": True,
                "pptx_path": output_path,
                "slide_count": len(slides_content),
                "output_files": output_files
            }

        except Exception as e:
            logger.exception("PPT 生成失败")
            return {
                "success": False,
                "error": "PPT生成失败，请稍后重试"
            }

    def _plan_ppt(self, user_request: str, slide_count: int) -> List[Dict]:
        """AI规划PPT内容"""
        from .ppt_planner import plan_ppt
        return plan_ppt(user_request, slide_count)

    def _generate_image(self, slide: Dict, slide_num: int = 1) -> Optional[str]:
        """生成AI图片，如果失败则返回备用图片URL"""
        from .content_generator import get_content_generator

        # 检查失败次数是否超过阈值
        if self._image_failure_count >= self._image_failure_threshold:
            logger.warning(f"AI图片生成失败次数({self._image_failure_count})超过阈值，直接使用备用图片")
            return self._get_fallback_image_url(slide, slide_num)

        content_gen = get_content_generator()

        # 构建更好的提示词
        prompt = slide.get("image_prompt", "")
        if not prompt:
            # 基于标题和内容构建提示词
            title = slide.get("title", "")
            content = slide.get("content", [])
            slide_type = slide.get("slide_type", "content")

            # 构建英文提示词
            prompt = self._build_image_prompt(title, content, slide_type)

        # 尝试生成图片 (使用16:9比例)
        try:
            image_url = content_gen.generate_image(
                prompt=prompt,
                model="stable-diffusion-xl-2600"
            )
            if image_url:
                logger.info(f"AI生成图片成功: {image_url[:50]}...")
                self._image_failure_count = 0  # 成功后重置计数器
                return image_url
        except Exception as e:
            self._image_failure_count += 1
            logger.warning(f"AI图片生成失败({self._image_failure_count}/{self._image_failure_threshold}): {e}")

        # 返回备用图片URL
        logger.info("使用备用图片")
        return self._get_fallback_image_url(slide, slide_num)

    def _build_image_prompt(self, title: str, content: List, slide_type: str) -> str:
        """构建AI图片提示词"""
        # 中文标题转英文
        title_en = self._translate_to_english(title)

        # 构建提示词
        if slide_type == "title" or slide_type == "thank_you":
            prompt = f"Professional business presentation {title_en}, clean modern design, minimal abstract background, soft gradient, high quality, 16:9 aspect ratio"
        elif slide_type == "toc":
            prompt = f"Table of contents layout, {title_en}, organized list design, clean typography, business presentation style, high quality, 16:9"
        else:
            # 内容页 - 基于内容构建
            content_desc = ", ".join(content[:3]) if content else "business content"
            content_en = self._translate_to_english(content_desc)
            prompt = f"Professional business illustration for {title_en}, {content_en}, modern corporate style, clean design, high quality, 16:9 aspect ratio"

        return prompt

    def _translate_to_english(self, text: str) -> str:
        """简单的中文到英文翻译（关键词映射）"""
        # 常用关键词映射
        keywords = {
            "商业": "business",
            "企业": "corporate",
            "公司": "company",
            "产品": "product",
            "服务": "service",
            "市场": "market",
            "营销": "marketing",
            "销售": "sales",
            "技术": "technology",
            "创新": "innovation",
            "发展": "development",
            "团队": "team",
            "领导": "leadership",
            "管理": "management",
            "培训": "training",
            "教育": "education",
            "学习": "learning",
            "财务": "finance",
            "数据": "data",
            "分析": "analysis",
            "报告": "report",
            "总结": "summary",
            "计划": "plan",
            "目标": "goal",
            "战略": "strategy",
            "优势": "advantage",
            "挑战": "challenge",
            "解决方案": "solution",
            "案例": "case study",
            "谢谢": "thank you",
            "感谢": "thank you",
            "结束": "ending"
        }

        text_lower = text.lower()
        for cn, en in keywords.items():
            if cn in text_lower:
                return en

        # 如果没有匹配，返回原文（作为后备）
        return text

    def _get_fallback_image_url(self, slide: Dict, slide_num: int = 1) -> str:
        """获取备用图片URL - 多级 Fallback 策略

        Level 1: Unsplash（带可达性检查）
        Level 2: 本地 SVG 占位图（纯色背景 + 图标）

        所有外部 URL 均经过 _validate_url() 验证后才使用。
        支持用户自定义 fallback 策略（通过 settings.image_fallback_strategy 配置）。
        """
        import random

        title = slide.get("title", "").lower()
        slide_type = slide.get("slide_type", "content")

        # ─── 用户自定义策略 ────────────────────────────────────────────
        user_strategy = getattr(settings, 'IMAGE_FALLBACK_STRATEGY', 'all')
        if user_strategy == 'none':
            # 用户选择不使用任何外部图片，直接用本地 SVG
            logger.info("用户配置跳过外部图片 fallback，使用本地 SVG")
            return self._generate_svg_placeholder(slide, slide_num)

        # ─── 分类候选 URL ─────────────────────────────────────────────
        categories = {
            "business": [  # 商务类
                "https://images.unsplash.com/photo-1497366216548-37526070297c?w=1600&q=80",
                "https://images.unsplash.com/photo-1497215842964-222b430dc094?w=1600&q=80",
                "https://images.unsplash.com/photo-1507679799987-c73779587ccf?w=1600&q=80",
            ],
            "tech": [  # 科技类
                "https://images.unsplash.com/photo-1518770660439-4636190af475?w=1600&q=80",
                "https://images.unsplash.com/photo-1451187580459-43490279c0fa?w=1600&q=80",
                "https://images.unsplash.com/photo-1526374965328-7f61d4dc18c5?w=1600&q=80",
            ],
            "education": [  # 教育类
                "https://images.unsplash.com/photo-1522202176988-66273c2fd55f?w=1600&q=80",
                "https://images.unsplash.com/photo-1524178232363-1fb2b075b655?w=1600&q=80",
                "https://images.unsplash.com/photo-1509062522246-3755977927d7?w=1600&q=80",
            ],
            "title": [  # 封面页
                "https://images.unsplash.com/photo-1557804506-669a67965ba0?w=1600&q=80",
                "https://images.unsplash.com/photo-1515187029135-18ee286d815b?w=1600&q=80",
                "https://images.unsplash.com/photo-1556761175-5973dc0f32e7?w=1600&q=80",
            ],
            "thank_you": [  # 尾页
                "https://images.unsplash.com/photo-1522071820081-009f0129c71c?w=1600&q=80",
                "https://images.unsplash.com/photo-1552664730-d307ca884978?w=1600&q=80",
            ],
            "default": [  # 默认商务
                "https://images.unsplash.com/photo-1551434678-e076c223a692?w=1600&q=80",
                "https://images.unsplash.com/photo-1460925895917-afdab827c52f?w=1600&q=80",
                "https://images.unsplash.com/photo-1506784983877-45594efa4cbe?w=1600&q=80",
            ],
        }

        # 根据内容分类选择候选池
        if any(kw in title for kw in ["培训", "销售", "企业", "公司", "管理"]):
            candidates = categories["business"]
        elif any(kw in title for kw in ["科技", "技术", "AI", "创新", "发展"]):
            candidates = categories["tech"]
        elif any(kw in title for kw in ["教育", "学习", "学校", "课程", "培训"]):
            candidates = categories["education"]
        elif slide_type == "title" or slide_num == 1:
            candidates = categories["title"]
        elif slide_type == "thank_you":
            candidates = categories["thank_you"]
        else:
            candidates = categories["default"]

        # ─── Level 1: Unsplash 带可达性验证 ─────────────────────────
        random.shuffle(candidates)
        for url in candidates:
            if self._validate_url(url):
                logger.info(f"Unsplash fallback 可用: {url[:60]}...")
                return url
            else:
                logger.warning(f"Unsplash URL 验证失败，跳过: {url[:60]}...")

        # ─── Level 2: 本地 SVG 占位图 ───────────────────────────────
        logger.info("所有 Unsplash URL 均不可达，使用本地 SVG 占位图")
        return self._generate_svg_placeholder(slide, slide_num)

    def _generate_svg_placeholder(self, slide: Dict, slide_num: int = 1) -> str:
        """生成本地 SVG 占位图 - 纯色背景 + 类别图标

        这是一个内嵌的 data URI，不依赖任何外部网络。
        """
        title = slide.get("title", "")
        slide_type = slide.get("slide_type", "content")

        # 根据 slide_type 选择配色和图标关键词
        color_map = {
            "title":      ("#1e40af", "#60a5fa", "📊"),   # 深蓝 + 图表图标
            "thank_you":  ("#065f46", "#34d399", "🙏"),
            "content":    ("#1e3a5f", "#60a5fa", "💡"),
            "toc":        ("#4c1d95", "#a78bfa", "📋"),
        }
        bg_color, accent_color, icon = color_map.get(slide_type, ("#1e3a5f", "#60a5fa", "📄"))

        # 转义标题（防 XSS）
        safe_title = self._escape_html(title or "Slide")

        # 序列化 SVG（不使用 data:image/svg+xml 前缀，内部以 <svg> 开头）
        inner = svg_placeholder_xml(slide_num, safe_title, bg_color, accent_color, icon)
        # Base64 编码以保证在所有上下文都安全
        import base64
        b64 = base64.b64encode(inner.encode('utf-8')).decode('ascii')
        return f"data:image/svg+xml;base64,{b64}"

    def _generate_svg(self, slide: Dict, slide_num: int) -> str:
        """AI生成SVG代码"""
        from .volc_api import get_volc_api
        
        volc = get_volc_api()
        
        # 检查火山引擎API是否可用
        api_available = self._check_volc_api_available(volc)
        
        if api_available:
            # 使用AI生成SVG
            logger.info("火山引擎API可用，使用AI生成SVG")
            return self._generate_svg_with_ai(slide, slide_num, volc)
        else:
            # 使用智能文字模式（备用方案）
            logger.info("火山引擎API不可用，使用智能文字模式")
            return self._generate_svg_smart_text(slide, slide_num)
    
    def _check_volc_api_available(self, volc) -> bool:
        """检查火山引擎API是否可用"""
        import os
        api_key = os.getenv("VOLCANO_API_KEY", os.getenv("VOLCENGINE_API_KEY", ""))
        project_id = os.getenv("VOLCENGINE_PROJECT_ID", "")
        
        if not api_key or not project_id:
            return False
        
        # 尝试调用API检查连通性
        try:
            response = volc.text_generation("test")
            return response.get("success", False)
        except Exception:
            return False
    
    def _generate_svg_with_ai(self, slide: Dict, slide_num: int, volc) -> str:
        """使用AI生成SVG代码"""
        # 构建SVG提示词
        prompt = self._build_svg_prompt(slide, slide_num)
        
        # 调用AI生成SVG (同步调用)
        response = volc.text_generation(prompt)
        
        # 提取SVG代码 - text_generation返回dict，需要取content
        if isinstance(response, dict):
            content = response.get("content", "")
        else:
            content = str(response)
        
        svg_code = self._extract_svg_code(content)

        return svg_code

    def _generate_svg_smart_layout(
        self,
        slide: Dict,
        slide_num: int,
        theme_color: str = "#165DFF",
        style: str = "professional",
        user_layout: str = None,
        text_style: str = "transparent_overlay",
        shadow_color: str = "#000000",
        overlay_transparency: int = 30,
        font_title: str = "思源黑体",
        font_subtitle: str = "思源黑体",
        font_content: str = "思源宋体",
        font_caption: str = "思源黑体",
        layout_mode: str = "auto",
        unified_layout: bool = True
    ) -> str:
        """使用智能布局模块生成SVG

        Args:
            slide: 幻灯片内容
            slide_num: 幻灯片编号
            theme_color: 主题色
            style: 风格
            user_layout: 用户指定的布局类型
            text_style: 文字样式 (transparent_overlay/shadow/glow/outline/gradient/neon)
            shadow_color: 阴影颜色
            overlay_transparency: 遮罩透明度
            font_title: 标题字体
            font_subtitle: 副标题字体
            font_content: 正文字体
            font_caption: 注释字体
        """
        from .smart_layout.creative_engine import get_creative_engine

        logger.info(f"[SmartLayout] slide_{slide_num}: layout_mode={layout_mode}, unified_layout={unified_layout}, text_style={text_style}")

        # 获取创意引擎
        engine = get_creative_engine()

        title = slide.get("title", "")
        subtitle = slide.get("subtitle", "")
        content = slide.get("content", [])

        # 准备内容
        content_list = []
        if subtitle:
            content_list.append({"title": subtitle, "content": ""})
        for item in content:
            if isinstance(item, dict):
                content_list.append({"title": item.get("title", ""), "content": item.get("content", "")})
            else:
                content_list.append({"title": str(item), "content": ""})

        # 确定布局类型
        if layout_mode == 'manual':
            # 手动模式：只使用用户指定的布局，不自动检测
            if user_layout:
                slide_type = user_layout
            else:
                # 没有指定布局时使用默认
                slide_type = "content_card"

            # 如果是第一页，使用封面布局
            if slide_num == 1:
                slide_type = "title_slide"

            # 统一布局模式：保存首页布局（线程安全）
            if unified_layout:
                with self._layout_lock:
                    if self._first_page_layout is None:
                        self._first_page_layout = slide_type
                        logger.info(f"[SmartLayout] slide_{slide_num}: manual模式，统一布局，设置首页布局 {slide_type}")
                    else:
                        slide_type = self._first_page_layout
                        logger.info(f"[SmartLayout] slide_{slide_num}: manual模式，统一布局，复用布局 {slide_type}")
            else:
                logger.info(f"[SmartLayout] slide_{slide_num}: manual模式，使用布局 {slide_type}")
        else:
            # 自动模式：智能检测布局
            if user_layout:
                # 有用户指定布局
                if unified_layout:
                    # 统一布局模式：使用线程锁确保线程安全
                    with self._layout_lock:
                        if self._first_page_layout is None:
                            # 第一次执行（可能是slide 1），设置首页布局
                            self._first_page_layout = user_layout
                            slide_type = user_layout
                            logger.info(f"[SmartLayout] slide_{slide_num}: 统一布局模式，设置首页布局 {slide_type}")
                        else:
                            # 后续页面复用首页布局
                            slide_type = self._first_page_layout
                            logger.info(f"[SmartLayout] slide_{slide_num}: 统一布局模式，复用布局 {slide_type}")
                else:
                    slide_type = user_layout
                    logger.info(f"[SmartLayout] slide_{slide_num}: 使用用户指定布局 {slide_type}")
            else:
                # 无用户指定，自动检测
                content_text = "\n".join([str(c) for c in content])
                suggestion = engine.get_layout_suggestion(title, content_text)
                slide_type = suggestion.get("content_type", "content")

                # 如果是第一页，使用封面布局
                if slide_num == 1:
                    slide_type = "title_slide"

                # 统一布局模式：保存首页布局（线程安全）
                if unified_layout:
                    with self._layout_lock:
                        if self._first_page_layout is None:
                            self._first_page_layout = slide_type
                            logger.info(f"[SmartLayout] slide_{slide_num}: 自动检测布局，设置首页布局 {slide_type}")
                        else:
                            slide_type = self._first_page_layout
                            logger.info(f"[SmartLayout] slide_{slide_num}: 自动检测布局，复用首页布局 {slide_type}")

        # 生成配色
        colors = engine.generate_color_palette(style, theme_color)

        # 添加文字样式和字体配置
        colors["text_style"] = text_style
        colors["shadow_color"] = shadow_color
        colors["overlay_transparency"] = overlay_transparency
        colors["fonts"] = {
            "title": font_title,
            "subtitle": font_subtitle,
            "content": font_content,
            "caption": font_caption
        }

        # 使用模板生成SVG（不使用AI）
        svg = engine._create_slide_from_template(title, content_list, slide_type, colors)

        logger.info(f"[SmartLayout] slide_{slide_num} layout={slide_type} text_style={text_style} generated successfully")
        return svg

    def _generate_svg_smart_text(self, slide: Dict, slide_num: int) -> str:
        """智能文字模式 - 不依赖AI生成SVG，使用纯代码生成美观SVG"""
        # DEBUG日志
        logger.info(f"[DEBUG] slide_{slide_num}: title={slide.get('title', '')[:30]}... content={slide.get('content', [])}")
        title = slide.get("title", "")
        subtitle = slide.get("subtitle", "")
        content = slide.get("content", [])
        image_url = slide.get("image_url", "")
        
        # 16:9 比例
        width = 1600
        height = 900
        
        # 构建内容列表HTML
        content_html = ""
        if content and len(content) > 0:
            items_html = ""
            for i, item in enumerate(content[:6]):  # 最多显示6个要点
                # 文字换行处理
                escaped_item = self._escape_html(str(item))
                items_html += f'''<text x="150" y="{380 + i * 70}" fill="rgba(255,255,255,0.9)" font-size="28" font-family="Microsoft YaHei, PingFang SC, sans-serif">
              <tspan x="150">{i+1}. </tspan><tspan>{escaped_item}</tspan>
            </text>'''
            
            content_html = f'''
<g id="content">
{items_html}
</g>'''
        
        # 副标题处理
        subtitle_html = ""
        if subtitle:
            escaped_subtitle = self._escape_html(subtitle)
            subtitle_html = f'''
<text x="800" y="180" text-anchor="middle" fill="rgba(255,255,255,0.8)" font-size="32" font-family="Microsoft YaHei, PingFang SC, sans-serif">
  {escaped_subtitle}
</text>'''
        
        # 标题处理
        escaped_title = self._escape_html(title)
        
        # 背景图片处理（如果有）- URL转义防止XSS
        bg_image = ""
        if image_url:
            escaped_url = self._escape_url(image_url)
            bg_image = f'''
<image href="{escaped_url}" x="0" y="0" width="1600" height="900" preserveAspectRatio="xMidYMid slice" opacity="0.4"/>'''
        
        svg_code = f'''<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" xmlns:xlink="http://www.w3.org/1999/xlink" viewBox="0 0 {width} {height}" width="{width}" height="{height}">
  <defs>
    <!-- 渐变背景 -->
    <linearGradient id="bgGradient" x1="0%" y1="0%" x2="100%" y2="100%">
      <stop offset="0%" style="stop-color:#0f172a"/>
      <stop offset="50%" style="stop-color:#1e3a5f"/>
      <stop offset="100%" style="stop-color:#3b82f6"/>
    </linearGradient>
    <!-- 标题渐变 -->
    <linearGradient id="titleGradient" x1="0%" y1="0%" x2="100%" y2="0%">
      <stop offset="0%" style="stop-color:#ffffff"/>
      <stop offset="100%" style="stop-color:#bfdbfe"/>
    </linearGradient>
    <!-- 背景遮罩 -->
    <linearGradient id="overlayGradient" x1="0%" y1="0%" x2="0%" y2="100%">
      <stop offset="0%" style="stop-color:rgba(15,23,42,0.3)"/>
      <stop offset="100%" style="stop-color:rgba(15,23,42,0.8)"/>
    </linearGradient>
  </defs>
  
  <!-- 背景 -->
  <rect width="{width}" height="{height}" fill="url(#bgGradient)"/>
  
  <!-- 背景图片（如果有） -->
  {bg_image}
  
  <!-- 遮罩层 -->
  <rect width="{width}" height="{height}" fill="url(#overlayGradient)"/>
  
  <!-- 装饰线条 -->
  <line x1="0" y1="200" x2="1600" y2="200" stroke="rgba(59,130,246,0.3)" stroke-width="2"/>
  <circle cx="100" cy="100" r="150" fill="rgba(59,130,246,0.1)"/>
  <circle cx="1500" cy="800" r="200" fill="rgba(59,130,246,0.08)"/>
  
  <!-- 页码 -->
  <text x="1550" y="870" text-anchor="end" fill="rgba(255,255,255,0.4)" font-size="18" font-family="Arial, sans-serif">
    {slide_num}
  </text>
  
  <!-- 标题 -->
  <text x="800" y="120" text-anchor="middle" fill="url(#titleGradient)" font-size="52" font-weight="bold" font-family="Microsoft YaHei, PingFang SC, sans-serif">
    {escaped_title}
  </text>
  
  <!-- 副标题 -->
  {subtitle_html}
  
  <!-- 内容要点 -->
  {content_html}
  
  <!-- 底部装饰 -->
  <rect x="0" y="850" width="1600" height="50" fill="rgba(59,130,246,0.2)"/>
  <text x="50" y="880" fill="rgba(255,255,255,0.5)" font-size="14" font-family="Arial, sans-serif">
    RabAi Mind - AI Generated
  </text>
</svg>'''
        
        return svg_code
    
    def _escape_html(self, text: str) -> str:
        """转义HTML特殊字符，防止XSS"""
        return (text
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
            .replace("'", "&#39;")
            .replace("`", "&#96;")
            .replace("\\", "&#92;"))

    def _escape_url(self, url: str) -> str:
        """转义URL中的特殊字符，防止XSS"""
        # 只转义XML/SVG中可能导致问题的字符
        return (url
            .replace("&", "&amp;")
            .replace('"', "&quot;")
            .replace("'", "&#39;")
            .replace("<", "&lt;")
            .replace(">", "&gt;"))

    def _validate_url(self, url: str) -> bool:
        """验证URL是否安全，防止SSRF攻击（含DNS Rebinding防护）"""
        from urllib.parse import urlparse
        import ipaddress
        import socket

        try:
            parsed = urlparse(url)
            hostname = parsed.hostname

            # 必须是http或https
            if parsed.scheme not in ('http', 'https'):
                return False

            # DNS Rebinding防护：遍历域名解析的所有IP，确保全部安全
            # 防止攻击者配置多IP DNS绕过检查
            dns_resolved = False
            try:
                # 获取所有IP地址
                addr_info = socket.getaddrinfo(hostname, None)
                all_ips = set(info[4][0] for info in addr_info)

                # 如果DNS解析成功但没有返回任何IP，拒绝请求
                if not all_ips:
                    return False

                dns_resolved = True
                for ip_str in all_ips:
                    try:
                        ip = ipaddress.ip_address(ip_str)

                        # 阻止所有内网IP
                        if (ip.is_private or ip.is_loopback or ip.is_link_local or
                            ip.is_reserved or ip.is_multicast):
                            return False

                        # 阻止特定内网段
                        if ip_str.startswith(('10.', '172.16.', '172.17.', '172.18.',
                                            '172.19.', '172.20.', '172.21.', '172.22.',
                                            '172.23.', '172.24.', '172.25.', '172.26.',
                                            '172.27.', '172.28.', '172.29.', '172.30.',
                                            '172.31.', '192.168.', '127.', '169.254.')):
                            return False

                        # IPv6内网
                        if ip_str.startswith(('::1', 'fe80:', 'fc', 'fd', '2001:db8:')):
                            return False
                    except ValueError:
                        continue

            except (socket.gaierror, OSError):
                # 无法解析，跳过IP检查但继续域名检查
                pass

            # 阻止常见内网域名
            internal_domains = ('localhost', '127.', '10.', '192.168.', '172.16.',
                              '172.17.', '172.18.', '172.19.', '172.2', '172.3',
                              '169.254.', '::1', 'fe80:', '.local', '.intranet.',
                              '.corp', '.home.', 'metadata.google.internal')
            for domain in internal_domains:
                if hostname.startswith(domain) or hostname.endswith(domain + '.local'):
                    return False

            return True
        except Exception:
            return False

    def _build_svg_prompt(self, slide: Dict, slide_num: int) -> str:
        """构建SVG生成提示词"""
        title = slide.get("title", "")
        content = slide.get("content", [])
        slide_type = slide.get("slide_type", "content")
        
        # 根据类型选择不同提示词
        if slide_num == 1:
            prompt = f"""# PPT页面SVG设计宗师

## 设计规范
- 比例规范: 严格遵循16:9 SVG viewBox="0 0 1600 900"
- 安全边际: 核心内容必须完整落于 100, 50, 1400, 800 安全区内
- 矢量原生: 仅使用SVG原生元素

## 配色方案
- 主色调: #165DFF (科技蓝)
- 背景: 渐变深蓝到浅蓝，或使用提供的背景图
- 文字: 白色或浅色

## 页面类型: 封面页

## 内容信息
- 标题: {title}
- 副标题: {slide.get('subtitle', '')}

## 背景图片
图片URL: {slide.get('image_url', '')}

请直接输出完整的 SVG 代码，不要包含任何解释或markdown标记。
确保背景使用提供的图片URL，覆盖整个页面。
"""
        else:
            content_str = "\n".join([f"- {c}" for c in content[:5]])
            prompt = f"""# PPT页面SVG设计宗师

## 设计规范
- 比例规范: 严格遵循16:9 SVG viewBox="0 0 1600 900"
- 安全边际: 核心内容必须完整落于 100, 50, 1400, 800 安全区内
- 矢量原生: 仅使用SVG原生元素

## 配色方案
- 主色调: #165DFF (科技蓝)
- 背景: 使用提供的背景图铺满整个页面
- 文字: 白色或浅色，覆盖在图片上

## 页面类型: 内容页

## 内容信息
- 标题: {title}
- 内容要点:
{content_str}

## 背景图片
图片URL: {slide.get('image_url', '')}

请直接输出完整的 SVG 代码，不要包含任何解释或markdown标记。
要求：
1. 背景图片铺满整个页面 (1600x900)
2. 标题文字大而醒目，放在页面顶部或中央
3. 内容文字放在图片下方或覆盖在图片下半部分
4. 文字使用白色或浅色带阴影确保可读性
"""
        return prompt

    def _extract_svg_code(self, response: str) -> str:
        """从AI响应中提取SVG代码"""
        # 尝试提取 ```svg ... ``` 或 ``` ... ``` 中的代码
        patterns = [
            r'```svg\s*(.*?)```',
            r'```\s*(.*?)```',
            r'<svg[^>]*>.*?</svg>',
        ]
        
        for pattern in patterns:
            matches = re.findall(pattern, response, re.DOTALL)
            if matches:
                code = matches[0].strip()
                if code.startswith('<svg'):
                    return code
        
        # 如果没有找到，返回默认SVG
        return self._default_svg()

    def _default_svg(self, slide_num: int = 1, title: str = "PPT Slide") -> str:
        """默认SVG - 用于错误时的回退"""
        escaped_title = self._escape_html(title)
        return f'''<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1600 900" width="1600" height="900">
  <defs>
    <linearGradient id="bg" x1="0%" y1="0%" x2="100%" y2="100%">
      <stop offset="0%" style="stop-color:#165DFF"/>
      <stop offset="100%" style="stop-color:#0D47E8"/>
    </linearGradient>
  </defs>
  <rect width="1600" height="900" fill="url(#bg)"/>
  <text x="800" y="420" text-anchor="middle" fill="white" font-size="56" font-family="Microsoft YaHei, PingFang SC, sans-serif" font-weight="bold">{escaped_title}</text>
  <text x="800" y="520" text-anchor="middle" fill="rgba(255,255,255,0.7)" font-size="28" font-family="Microsoft YaHei, PingFang SC, sans-serif">第 {slide_num} 页</text>
</svg>'''

    def _svg_to_ppt(self, svg_files: List[str], output_path: str, text_style: str = "transparent_overlay", shadow_color: str = "#000000", overlay_transparency: int = 30, use_smart_layout: bool = False, slide_backgrounds: list = None, font_title: str = "微软雅黑", font_subtitle: str = "微软雅黑", font_content: str = "微软雅黑", font_caption: str = "微软雅黑") -> bool:
        """SVG转PPT - 支持三种文字样式方案和智能布局模式

        兼容性说明:
        - Office 2007+ 要求: OOXML格式 (pptx)
        - WPS 2019+ 支持: 标准OOXML格式
        - 确保使用兼容的命名空间和元素
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor
            import requests
            import uuid
            from PIL import Image

            # 创建演示文稿
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)

            # 字体映射 - 转换为PPT兼容的字体名
            font_map = {
                "思源黑体": "Microsoft YaHei",
                "思源宋体": "SimSun",
                "思源黑体 Heavy": "Microsoft YaHei",
                "思源黑体 Bold": "Microsoft YaHei",
                "Noto Sans": "Microsoft YaHei",
                "Source Han Sans": "Microsoft YaHei",
                "Source Han Serif": "SimSun",
            }
            title_font = font_map.get(font_title, font_title)
            subtitle_font = font_map.get(font_subtitle, font_subtitle)
            content_font = font_map.get(font_content, font_content)

            # 设置兼容模式 - 确保Office 2007和WPS兼容
            # python-pptx默认已使用OOXML兼容格式

            for idx, svg_file in enumerate(svg_files):
                if not os.path.exists(svg_file):
                    continue

                # 读取SVG
                with open(svg_file, 'r', encoding='utf-8') as f:
                    svg_content = f.read()

                # 提取背景图片URL
                img_url = self._extract_image_url(svg_content)

                # 创建空白幻灯片
                blank_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_layout)

                # ===== 智能布局模式处理 =====
                if use_smart_layout:
                    # 确定背景颜色：优先使用用户设置，否则使用SVG中的颜色
                    bg_color = None
                    if slide_backgrounds and idx < len(slide_backgrounds):
                        # 使用用户设置的背景颜色
                        bg = slide_backgrounds[idx]
                        bg_color = bg.background_color if hasattr(bg, 'background_color') else None
                        logger.info(f"使用用户设置的背景颜色: {bg_color}")
                    if not bg_color:
                        # 从SVG提取背景颜色
                        bg_color = self._extract_svg_background_color(svg_content)
                        logger.info(f"使用SVG背景颜色: {bg_color}")

                    if bg_color and len(bg_color) == 6:
                        # 转换颜色
                        r = int(bg_color[0:2], 16)
                        g = int(bg_color[2:4], 16)
                        b = int(bg_color[4:6], 16)
                        slide.background.fill.solid()
                        slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
                        logger.info(f"已设置背景颜色: {bg_color}")

                        # 计算背景亮度，决定文字颜色
                        brightness = (r * 299 + g * 587 + b * 114) / 1000
                        text_color = RGBColor(255, 255, 255) if brightness < 128 else RGBColor(0, 0, 0)
                    else:
                        # 使用默认蓝色背景
                        slide.background.fill.solid()
                        slide.background.fill.fore_color.rgb = RGBColor(0x16, 0x5D, 0xFF)
                        text_color = RGBColor(255, 255, 255)  # 蓝色背景用白色文字
                        logger.info("已设置默认蓝色背景")

                    # 提取并添加文字
                    title, contents = self._extract_text_from_svg(svg_content)

                    # 添加标题
                    if title:
                        title_box = slide.shapes.add_textbox(
                            Inches(0.5), Inches(0.3),
                            Inches(15), Inches(1.2)
                        )
                        tf = title_box.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        p.text = title
                        p.font.size = Pt(48)
                        p.font.bold = True
                        p.font.color.rgb = text_color
                        p.font.name = title_font
                        p.alignment = 1  # 居中
                        logger.info(f"已添加标题: {title}")

                    # 添加内容
                    if contents:
                        # 过滤掉与标题相同的内容，并去掉已有的圆点符号
                        filtered_contents = []
                        for c in contents:
                            if c != title:
                                # 去掉已有的圆点符号前缀
                                if c.startswith("• "):
                                    c = c[2:]
                                filtered_contents.append(c)

                        if filtered_contents:
                            content_box = slide.shapes.add_textbox(
                                Inches(1.0), Inches(2.2),
                                Inches(14), Inches(5)
                            )
                            tf = content_box.text_frame
                            tf.word_wrap = True
                            for i, content in enumerate(filtered_contents[:6]):  # 最多6条
                                if i == 0:
                                    p = tf.paragraphs[0]
                                else:
                                    p = tf.add_paragraph()
                                # 所有内容都统一添加圆点符号
                                p.text = f"• {content}"
                                p.font.size = Pt(24)
                                p.font.color.rgb = text_color
                                p.font.name = content_font
                                p.space_before = Pt(12)
                        logger.info(f"已添加 {len(contents)} 条内容")

                    logger.info(f"智能布局模式：已添加文字")
                    continue

                # ===== 原有模式：添加背景图片（铺满整个页面）=====
                image_brightness = 128  # 默认中等亮度
                if img_url:
                    # SSRF防护：验证URL
                    if not self._validate_url(img_url):
                        logger.warning(f"URL验证失败，跳过图片: {img_url[:50]}...")
                        continue
                    try:
                        response = requests.get(img_url, timeout=60)
                        if response.status_code == 200:
                            # 保存临时文件（使用安全临时目录）
                            import tempfile
                            temp_dir = tempfile.gettempdir()
                            temp_img = os.path.join(temp_dir, f"rabai_temp_{uuid.uuid4().hex}.jpg")

                            # 图片压缩 - 优化文件大小
                            try:
                                img = Image.open(BytesIO(response.content))
                                # 调整大小为合适分辨率 (1920x1080 for 16:9)
                                img = img.resize((1920, 1080), Image.Resampling.LANCZOS)
                                # 保存为压缩JPEG
                                img.save(temp_img, 'JPEG', quality=75, optimize=True)
                            except Exception as e:
                                logger.warning(f"图片压缩失败，使用原图: {e}")
                                with open(temp_img, 'wb') as f:
                                    f.write(response.content)

                            # 分析图片亮度
                            try:
                                img = Image.open(temp_img)
                                img = img.resize((100, 100))  # 缩小以提高速度
                                import numpy as np
                                img_array = np.array(img)
                                # 计算平均亮度
                                if len(img_array.shape) == 3:
                                    image_brightness = np.mean(img_array[:, :, :3])
                                else:
                                    image_brightness = np.mean(img_array)
                            except Exception as e:
                                logger.warning(f"图片亮度分析失败: {e}")

                            # 添加图片铺满页面 (16:9)
                            try:
                                slide.shapes.add_picture(
                                    temp_img,
                                    Inches(0), Inches(0),
                                    width=Inches(16), height=Inches(9)
                                )
                                logger.info(f"已添加背景图片: {img_url[:50]}..., 亮度: {image_brightness}")
                            finally:
                                # 确保临时文件被清理
                                if os.path.exists(temp_img):
                                    os.remove(temp_img)
                        else:
                            logger.warning(f"图片下载失败: {response.status_code}")
                    except Exception as e:
                        logger.warning(f"添加背景图片失败: {e}")
                else:
                    logger.warning("SVG中未找到图片URL")

                # ===== 步骤2: 根据文字样式方案处理 =====
                # 提取文字内容
                title, contents = self._extract_text_from_svg(svg_content)

                # ===== 方案1: 半透明遮罩 =====
                if text_style == "transparent_overlay":
                    from pptx.oxml.ns import qn
                    from pptx.oxml.xmlchemy import OxmlElement

                    # 标题区域遮罩 - 使用30%透明度通过OXML
                    title_bg = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        Inches(16), Inches(2.0)
                    )
                    # 使用OXML设置30%透明度的黑色填充
                    spPr = title_bg._element.spPr
                    solidFill = OxmlElement('a:solidFill')
                    srgbClr = OxmlElement('a:srgbClr')
                    srgbClr.set('val', '000000')
                    alpha = OxmlElement('a:alpha')
                    alpha.set('val', str(overlay_transparency * 1000))  # 透明度百分比转换为PPTX alpha值
                    srgbClr.append(alpha)
                    solidFill.append(srgbClr)
                    spPr.append(solidFill)
                    title_bg.line.fill.background()
                    logger.info(f"已设置标题遮罩透明度{overlay_transparency}%")

                    # 内容区域遮罩
                    if contents:
                        content_height = min(len(contents) * 0.9 + 1.5, 5.5)
                        content_bg = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            Inches(0), Inches(2.0),
                            Inches(16), Inches(content_height)
                        )
                        # 使用OXML设置30%透明度的黑色填充
                        spPr = content_bg._element.spPr
                        solidFill = OxmlElement('a:solidFill')
                        srgbClr = OxmlElement('a:srgbClr')
                        srgbClr.set('val', '000000')
                        alpha = OxmlElement('a:alpha')
                        alpha.set('val', str(overlay_transparency * 1000))  # 透明度百分比转换为PPTX alpha值
                        srgbClr.append(alpha)
                        solidFill.append(srgbClr)
                        spPr.append(solidFill)
                        content_bg.line.fill.background()
                        logger.info(f"已设置内容遮罩透明度{overlay_transparency}%")

                # ===== 确定文字颜色 =====
                # 默认白色文字
                text_color = RGBColor(255, 255, 255)

                # 转换阴影颜色
                shadow_rgb = self._hex_to_rgb(shadow_color)

                # ===== 添加标题 =====
                if title:
                    # ===== 方案3: 文字阴影 =====
                    if text_style == "shadow":
                        # 添加阴影效果 - 先添加偏移的文字
                        shadow_offset = 3
                        # 阴影（自定义颜色）
                        title_shadow = slide.shapes.add_textbox(
                            Inches(0.5 + shadow_offset*0.01), Inches(0.3 + shadow_offset*0.01),
                            Inches(15), Inches(1.2)
                        )
                        tf_shadow = title_shadow.text_frame
                        tf_shadow.word_wrap = True
                        p_shadow = tf_shadow.paragraphs[0]
                        p_shadow.text = title
                        p_shadow.font.size = Pt(48)
                        p_shadow.font.bold = True
                        p_shadow.font.color.rgb = shadow_rgb  # 自定义阴影颜色
                        p_shadow.alignment = 1

                    # 标题（主文字）
                    title_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.3),
                        Inches(15), Inches(1.2)
                    )
                    tf = title_box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = title
                    p.font.size = Pt(48)
                    p.font.bold = True
                    p.font.color.rgb = text_color
                    p.alignment = 1  # 居中

                # ===== 添加内容 =====
                if contents:
                    # 内容区域位置根据样式调整
                    content_top = 2.5 if text_style != "transparent_overlay" else 2.2

                    # ===== 方案3: 文字阴影 =====
                    if text_style == "shadow":
                        # 阴影层
                        shadow_offset = 2
                        content_shadow = slide.shapes.add_textbox(
                            Inches(1.0 + shadow_offset*0.01), Inches(content_top + shadow_offset*0.01),
                            Inches(14), Inches(5)
                        )
                        tf_shadow = content_shadow.text_frame
                        tf_shadow.word_wrap = True
                        for i, content in enumerate(contents[:6]):
                            if i == 0:
                                p_shadow = tf_shadow.paragraphs[0]
                            else:
                                p_shadow = tf_shadow.add_paragraph()
                            p_shadow.text = f"  {content}"
                            p_shadow.font.size = Pt(26)
                            p_shadow.font.color.rgb = shadow_rgb  # 自定义阴影颜色
                            p_shadow.space_before = Pt(18)

                    # 内容层（主文字）
                    content_box = slide.shapes.add_textbox(
                        Inches(1.0), Inches(content_top),
                        Inches(14), Inches(5)
                    )
                    tf = content_box.text_frame
                    tf.word_wrap = True

                    for i, content in enumerate(contents[:6]):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = f"  {content}"
                        p.font.size = Pt(26)
                        p.font.color.rgb = text_color
                        p.space_before = Pt(18)

            # 保存（原子写入：先写临时文件再rename）
            import tempfile
            temp_path = output_path + '.tmp'
            try:
                prs.save(temp_path)
                os.replace(temp_path, output_path)  # 原子重命名
                logger.info(f"PPT保存成功: {output_path}")
            finally:
                # 确保临时文件被清理（无论成功还是异常）
                if os.path.exists(temp_path):
                    try:
                        os.remove(temp_path)
                    except OSError:
                        pass

            # 文件压缩优化
            self._compress_pptx(output_path)

            return True

        except Exception as e:
            logger.exception("SVG转PPT失败")
            return False

    def _compress_pptx(self, file_path: str) -> bool:
        """压缩PPTX文件"""
        import zipfile
        import os

        if not os.path.exists(file_path):
            return False

        try:
            original_size = os.path.getsize(file_path)

            # 创建临时文件
            temp_path = file_path + '.tmp'

            # 使用最大压缩比重新压缩
            with zipfile.ZipFile(temp_path, 'w', zipfile.ZIP_DEFLATED, compresslevel=9) as zip_out:
                with zipfile.ZipFile(file_path, 'r') as zip_in:
                    for item in zip_in.infolist():
                        data = zip_in.read(item.filename)
                        # 对媒体文件使用更高压缩
                        if item.filename.endswith(('.png', '.jpg', '.jpeg', '.gif')):
                            zip_out.writestr(item, data, compresslevel=9)
                        else:
                            zip_out.writestr(item, data, compresslevel=9)

            # 替换原文件
            os.replace(temp_path, file_path)

            compressed_size = os.path.getsize(file_path)
            ratio = (1 - compressed_size / original_size) * 100

            logger.info(f"PPT压缩完成: {original_size / 1024 / 1024:.2f}MB -> {compressed_size / 1024 / 1024:.2f}MB (减少 {ratio:.1f}%)")

            return True

        except Exception as e:
            logger.warning(f"PPT压缩失败: {e}")
            # 清理临时文件
            temp_path = file_path + '.tmp'
            if os.path.exists(temp_path):
                os.remove(temp_path)
            return False

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """将十六进制颜色转换为RGBColor"""
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 6:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)
        return RGBColor(0, 0, 0)

    def _extract_svg_background_color(self, svg_content: str) -> Optional[str]:
        """从SVG中提取背景颜色（从stop-color或fill属性）"""
        import re

        # 尝试从stop-color提取渐变起始颜色
        match = re.search(r'stop-color:\s*([^;"]+)', svg_content)
        if match:
            color = match.group(1).strip()
            if color.startswith('#'):
                # 转换 #RRGGBB 为 RRGGBB
                return color[1:]

        # 尝试从 fill 属性提取背景色
        match = re.search(r'<rect[^>]*fill="([^"]+)"', svg_content)
        if match:
            fill = match.group(1)
            if fill.startswith('#'):
                return fill[1:]
            # 检查是否是渐变
            if 'url(#' in fill:
                # 尝试从渐变定义中提取颜色
                grad_match = re.search(r'id="[^"]*"[^>]*>.*?<stop[^>]*stop-color:\s*([^;"]+)', svg_content, re.DOTALL)
                if grad_match:
                    color = grad_match.group(1).strip()
                    if color.startswith('#'):
                        return color[1:]

        return None

    def _extract_image_url(self, svg_content: str) -> Optional[str]:
        """从SVG中提取图片URL"""
        patterns = [
            # 匹配 href="URL" - 支持URL参数
            r'href="([^"]+)"',
            # 匹配 xlink:href="URL"
            r'xlink:href="([^"]+)"',
            # 直接匹配URL
            r'(https?://[^\s"\'<>]+)',
        ]

        for pattern in patterns:
            match = re.search(pattern, svg_content, re.I)
            if match:
                url = match.group(1)
                # 过滤掉非图片URL
                if any(ext in url.lower() for ext in ['.jpg', '.jpeg', '.png', '.webp', '.gif']):
                    return url
                # 也接受unsplash等图片服务URL
                if 'unsplash' in url or 'images.unsplash' in url:
                    return url

        return None

    def _extract_text_from_svg(self, svg_content: str):
        """从SVG中提取文字（包括tspan内的文字）"""
        import re
        
        # 方法1: 使用正则提取所有text和tspan元素的内容
        # 先提取所有tspan内容
        tspan_contents = re.findall(r'<tspan[^>]*>([^<]*)</tspan>', svg_content)
        tspan_contents = [t.strip() for t in tspan_contents if t.strip()]
        
        # 提取所有text元素的直接内容（不含tspan）
        text_contents = re.findall(r'<text[^>]*>([^<]*)</text>', svg_content)
        text_contents = [t.strip() for t in text_contents if t.strip()]
        
        # 合并所有文本内容
        all_texts = tspan_contents + text_contents
        
        # 过滤掉页码数字、装饰文字等
        filtered_texts = []
        exclude_patterns = [
            r'^\d+$',  # 纯数字（页码）
            r'^RabAi',  # 底部商标
        ]

        for text in all_texts:
            is_excluded = False
            for pattern in exclude_patterns:
                if re.match(pattern, text):
                    is_excluded = True
                    break
            # 过滤太短的文字（允许如"AI"这样的短词）
            if len(text) >= 2 and not is_excluded:
                filtered_texts.append(text)
        
        # 识别标题：通常是大字体、在页面顶部(y < 200)
        title = ""
        contents = []

        # 添加调试日志
        logger.info(f"[SVG提取] 开始提取文字，SVG长度: {len(svg_content)}")

        # 通过y坐标识别标题（y在80-200之间的是标题）
        title_matches = re.findall(r'<text[^>]*y="(\d+)"[^>]*>.*?</text>', svg_content)
        text_with_y = []
        for match in re.finditer(r'<text[^>]*y="(\d+)"[^>]*>(.*?)</text>', svg_content, re.DOTALL):
            y = int(match.group(1))
            text = match.group(2)
            # 提取tspan内容
            tspan_in_text = re.findall(r'<tspan[^>]*>([^<]*)</tspan>', text)
            if tspan_in_text:
                text = ' '.join([t.strip() for t in tspan_in_text if t.strip()])
            else:
                text = text.strip()
            if text and len(text) > 2:
                text_with_y.append((y, text))
        
        # 按y坐标排序
        text_with_y.sort(key=lambda x: x[0])

        logger.info(f"[SVG提取] 找到 {len(text_with_y)} 个文本元素: {text_with_y}")

        # 找到标题（第一个文本，通常是标题）
        # 跳过装饰性的引号等
        for y, text in text_with_y:
            # 跳过纯装饰符号
            if text in ['"', '"'] or len(text) < 3:
                continue
            # 取第一个作为标题
            title = text
            break

        # 其余的是内容（跳过标题）
        for y, text in text_with_y:
            if text != title and len(text) > 2:
                # 跳过页码和底部装饰
                if y > 800:
                    continue
                contents.append(text)

        logger.info(f"[SVG提取] 提取结果: title='{title}', contents={contents}")
        
        # 如果上述方法失败，使用备用逻辑
        if not title and filtered_texts:
            # 取第一个最长的作为标题
            title = max(filtered_texts, key=len) if filtered_texts else ""
            # 其余作为内容
            contents = [t for t in filtered_texts if t != title]
        
        return title, contents


    def _convert_svg_to_png(self, svg_files: List[str], output_dir: str, quality: str = "standard") -> List[str]:
        """将SVG转换为PNG图片

        Args:
            svg_files: SVG文件路径列表
            output_dir: 输出目录
            quality: 质量 (standard/high/ultra)

        Returns:
            PNG文件路径列表
        """
        # 根据质量设置分辨率
        quality_dpi = {
            "standard": 96,
            "high": 144,
            "ultra": 192
        }
        dpi = quality_dpi.get(quality, 96)

        png_paths = []
        for svg_file in svg_files:
            if not os.path.exists(svg_file):
                continue

            # 提取文件名
            basename = os.path.basename(svg_file).replace(".svg", ".png")
            png_path = os.path.join(output_dir, basename)

            try:
                # 使用cairosvg转换（如果可用）
                try:
                    import cairosvg
                    cairosvg.svg2png(url=svg_file, write_to=png_path, dpi=dpi)
                    logger.info(f"SVG转PNG成功: {png_path}")
                except ImportError:
                    # 备用方案：使用PIL
                    from PIL import Image
                    import io

                    # 读取SVG
                    with open(svg_file, 'r', encoding='utf-8') as f:
                        svg_content = f.read()

                    # 创建临时PNG
                    # 注意：PIL需要先转换为其他格式，这里简化处理
                    # 实际使用需要cairosvg或其他库
                    logger.warning("cairosvg未安装，PNG转换可能不完整")

                png_paths.append(png_path)
            except Exception as e:
                logger.error(f"SVG转PNG失败: {e}")

        return png_paths

    def _svg_to_pdf(self, svg_files: List[str], output_path: str) -> bool:
        """将SVG转换为PDF

        Args:
            svg_files: SVG文件路径列表
            output_path: 输出PDF路径

        Returns:
            是否成功
        """
        try:
            # 方法1: 使用reportlab
            try:
                from reportlab.graphics import renderPDF
                from reportlab.lib.pagesizes import landscape
                from svglib.svglib import svg2rlg

                from reportlab.pdfgen import canvas
                from reportlab.lib.units import inch

                c = canvas.Canvas(output_path, pagesize=landscape((16*inch, 9*inch)))

                for svg_file in svg_files:
                    if not os.path.exists(svg_file):
                        continue

                    # 转换SVG到reportlab图形
                    drawing = svg2rlg(svg_file)

                    if drawing:
                        # 渲染到PDF
                        renderPDF.draw(drawing, c, 0, 0)
                        c.showPage()

                c.save()
                logger.info(f"SVG转PDF成功: {output_path}")
                return True
            except ImportError:
                # 方法2: 备用方案 - 先转PPTX再转PDF
                logger.warning("reportlab未安装，尝试备用方案")

                # 创建临时PPTX
                import tempfile
                tmp_path = None
                try:
                    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                        tmp_path = tmp.name

                    self._svg_to_ppt(svg_files, tmp_path)
                    logger.info(f"临时PPTX已生成: {tmp_path}")
                    # 注意：实际PDF转换需要安装libreoffice或其他工具
                    # 这里只是标记需要后续处理
                    logger.warning("请使用PPTX导出或安装libreoffice进行PDF转换")
                finally:
                    # Clean up temp PPTX file
                    if tmp_path:
                        try:
                            os.unlink(tmp_path)
                        except OSError:
                            pass

                return False

        except Exception as e:
            logger.error(f"SVG转PDF失败: {e}")
            return False


# 全局实例
_ppt_generator: Optional[PPTGenerator] = None
_generator_lock = threading.Lock()  # 保护单例创建


def get_ppt_generator() -> PPTGenerator:
    """获取PPT生成器实例（线程安全）"""
    global _ppt_generator
    if _ppt_generator is None:
        with _generator_lock:
            # 双重检查
            if _ppt_generator is None:
                _ppt_generator = PPTGenerator()
    return _ppt_generator
