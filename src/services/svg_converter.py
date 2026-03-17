# -*- coding: utf-8 -*-
"""
SVG to PPTX 转换模块 (简化版)

不依赖 cairo 库，使用 python-pptx 直接创建 PPT
"""

import os
import uuid
import logging
import re
import requests
from typing import Optional, Union, Tuple, List
from io import BytesIO

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm, Emu
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available, PPTX creation will not work")


def _download_image(url: str) -> Optional[BytesIO]:
    """下载图片并返回 BytesIO"""
    try:
        response = requests.get(url, timeout=30)
        if response.status_code == 200:
            return BytesIO(response.content)
    except Exception as e:
        logger.warning(f"Failed to download image: {e}")
    return None


def create_pptx_from_svgs(
    svg_files: List[str],
    output_path: str,
    slide_width: float = 16.0,
    slide_height: float = 9.0
) -> Tuple[bool, str]:
    """
    从多个 SVG 文件创建 PPTX，支持图片

    Args:
        svg_files: SVG 文件路径列表
        output_path: 输出 PPTX 路径
        slide_width: 幻灯片宽度（英寸）
        slide_height: 幻灯片高度（英寸）

    Returns:
        (成功标志, 错误信息)
    """
    if not PPTX_AVAILABLE:
        return False, "python-pptx not available"

    try:
        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(slide_width)
        prs.slide_height = Inches(slide_height)

        # 为每个 SVG 创建一个幻灯片
        for svg_path in svg_files:
            if not os.path.exists(svg_path):
                logger.warning(f"SVG not found: {svg_path}")
                continue

            # 尝试从 SVG 提取文本内容和图片
            title, content_lines, image_url = _extract_text_from_svg(svg_path)

            # 创建空白幻灯片
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            # 添加图片 (如果存在)
            if image_url:
                try:
                    img_io = _download_image(image_url)
                    if img_io:
                        # 保存为临时文件
                        temp_img = os.path.join("/tmp", f"temp_{uuid.uuid4().hex}.png")
                        with open(temp_img, 'wb') as f:
                            f.write(img_io.getvalue())

                        # 添加图片到右侧
                        slide.shapes.add_picture(
                            temp_img,
                            Inches(8.5), Inches(1.5),
                            width=Inches(7)
                        )
                        os.remove(temp_img)
                        logger.info(f"Added image to slide: {image_url[:50]}...")
                except Exception as e:
                    logger.warning(f"Failed to add image: {e}")

            # 添加标题
            if title:
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3),
                    Inches(7.5), Inches(1)
                )
                text_frame = title_box.text_frame
                text_frame.text = title
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(36)
                    paragraph.font.bold = True

            # 添加内容 (在图片左侧)
            if content_lines:
                content_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.5),
                    Inches(7.5), Inches(7)
                )
                text_frame = content_box.text_frame
                text_frame.word_wrap = True

                for i, line in enumerate(content_lines[:8]):  # 最多显示8行
                    if i == 0:
                        text_frame.text = line[:100] if len(line) > 100 else line
                    else:
                        p = text_frame.add_paragraph()
                        p.text = line[:100] if len(line) > 100 else line
                        p.font.size = Pt(18)

        # 保存文件
        prs.save(output_path)
        logger.info(f"Created PPTX: {output_path}")
        return True, ""

    except Exception as e:
        logger.error(f"Error creating PPTX: {e}")
        import traceback
        traceback.print_exc()
        return False, str(e)


def _extract_text_from_svg(svg_path: str) -> Tuple[Optional[str], List[str], Optional[str]]:
    """
    从 SVG 文件提取文本内容和图片URL

    Returns:
        (标题, 内容行列表, 图片URL)
    """
    try:
        with open(svg_path, 'r', encoding='utf-8') as f:
            content = f.read()

        title = None
        content_lines = []
        image_url = None

        # 提取标题 (第一个 text 元素)
        title_match = re.search(r'<text[^>]*>([^<]+)</text>', content)
        if title_match:
            title = title_match.group(1).strip()

        # 提取所有文本内容
        text_matches = re.findall(r'<text[^>]*>([^<]+)</text>', content)
        for i, text in enumerate(text_matches):
            text = text.strip()
            if text and text != title:
                # 过滤掉一些非内容文本
                if 'RabAi' not in text and 'PPT' not in text:
                    content_lines.append(text)

        # 提取图片 URL
        image_match = re.search(r'<image[^>]*href="([^"]+)"', content)
        if image_match:
            image_url = image_match.group(1)

        return title, content_lines[:10], image_url

    except Exception as e:
        logger.error(f"Error extracting text from SVG: {e}")
        return None, [], None


def get_pptx_slide_count(pptx_path: str) -> Tuple[int, str]:
    """
    获取 PPTX 幻灯片数量

    Args:
        pptx_path: PPTX 文件路径

    Returns:
        (数量, 错误信息)
    """
    if not PPTX_AVAILABLE:
        return 0, "python-pptx not available"

    try:
        if not os.path.exists(pptx_path):
            return 0, f"File not found: {pptx_path}"

        prs = Presentation(pptx_path)
        return len(prs.slides), ""

    except Exception as e:
        return 0, str(e)


def insert_svg_to_pptx(
    pptx_path: str,
    svg_path: str,
    slide_number: int = 1,
    x: Optional[Union[float, int]] = None,
    y: Optional[Union[float, int]] = None,
    width: Optional[Union[float, int]] = None,
    height: Optional[Union[float, int]] = None,
    output_path: Optional[str] = None,
    create_if_not_exists: bool = True
) -> Tuple[bool, str]:
    """
    将 SVG 内容插入到 PPTX 文件

    注意: 由于没有 cairo 库，这里会提取文本并添加到 PPT
    """
    if not PPTX_AVAILABLE:
        return False, "python-pptx not available"

    try:
        # 提取文本
        title, content_lines = _extract_text_from_svg(svg_path)

        # 打开或创建 PPTX
        if os.path.exists(pptx_path):
            prs = Presentation(pptx_path)
        else:
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)

        # 确保有足够的幻灯片
        while len(prs.slides) < slide_number:
            prs.slides.add_slide(prs.slide_layouts[6])

        # 获取幻灯片
        slide = prs.slides[slide_number - 1]

        # 添加标题
        if title:
            # 检查是否已有标题
            has_title = any(shape.name == "Title" for shape in slide.shapes if hasattr(shape, "name"))

            if not has_title:
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3),
                    Inches(15), Inches(1)
                )
                text_frame = title_box.text_frame
                text_frame.text = title
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(32)
                    paragraph.font.bold = True

        # 添加内容
        if content_lines:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8),
                Inches(15), Inches(6)
            )
            text_frame = content_box.text_frame
            text_frame.word_wrap = True

            for i, line in enumerate(content_lines[:8]):
                if i == 0:
                    text_frame.text = line
                else:
                    p = text_frame.add_paragraph()
                    p.text = line
                    p.font.size = Pt(18)

        # 保存
        output_path = output_path or pptx_path
        prs.save(output_path)

        return True, ""

    except Exception as e:
        import traceback
        traceback.print_exc()
        return False, str(e)
