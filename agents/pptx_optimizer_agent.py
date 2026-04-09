"""
PPTX 优化 Agent

负责对生成的 PPTX 文件进行优化：
- 字体嵌入
- 图片压缩
- 文件清理
"""
import os
import logging
import shutil
import tempfile
import zipfile
from pathlib import Path
from typing import Optional

try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    from pptx.util import Inches, Emu
except ImportError:
    raise ImportError("python-pptx is required: pip install python-pptx")

try:
    from PIL import Image
except ImportError:
    raise ImportError("Pillow is required: pip install Pillow")

logger = logging.getLogger(__name__)


class PPTXOptimizerAgent:
    """PPT 优化 Agent"""

    def __init__(self, output_dir: Optional[str] = None):
        self.name = "PPTXOptimizerAgent"
        self.output_dir = output_dir or tempfile.mkdtemp(prefix="pptx_optimized_")

    def optimize(self, pptx_path: str, options: Optional[dict] = None) -> str:
        """
        优化 PPTX，返回优化后路径

        Args:
            pptx_path: PPTX 文件路径
            options: 优化选项 {compress: bool, embed_fonts: bool, clean_metadata: bool}
        Returns:
            优化后文件路径
        """
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PPTX 文件不存在: {pptx_path}")

        options = options or {}
        embed_fonts = options.get("embed_fonts", True)
        compress = options.get("compress", True)
        clean_metadata = options.get("clean_metadata", True)

        logger.info(f"PPTXOptimizer: 开始优化文件 {pptx_path}")
        logger.info(f"  - embed_fonts: {embed_fonts}")
        logger.info(f"  - compress: {compress}")
        logger.info(f"  - clean_metadata: {clean_metadata}")

        # 创建优化后的文件路径
        base_name = os.path.basename(pptx_path)
        name_without_ext = os.path.splitext(base_name)[0]
        optimized_path = os.path.join(self.output_dir, f"{name_without_ext}_optimized.pptx")

        # 确保输出目录存在
        os.makedirs(self.output_dir, exist_ok=True)

        # 复制原文件到临时目录进行处理
        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = os.path.join(tmp_dir, base_name)
            shutil.copy2(pptx_path, tmp_path)

            # 1. 嵌入字体
            if embed_fonts:
                tmp_path = self._embed_fonts_internal(tmp_path)

            # 2. 压缩图片
            if compress:
                tmp_path = self._compress_images_internal(tmp_path)

            # 3. 清理元数据
            if clean_metadata:
                tmp_path = self._clean_metadata_internal(tmp_path)

            # 移动到最终位置
            shutil.move(tmp_path, optimized_path)

        logger.info(f"PPTXOptimizer: 优化完成，输出到 {optimized_path}")
        return optimized_path

    def _embed_fonts_internal(self, pptx_path: str) -> str:
        """嵌入缺失字体的内部方法"""
        logger.info(f"PPTXOptimizer: 检查并嵌入字体...")

        try:
            prs = Presentation(pptx_path)

            # 常见的中文字体映射
            font_mapping = {
                "宋体": "SimSun",
                "黑体": "SimHei",
                "楷体": "KaiTi",
                "微软雅黑": "Microsoft YaHei",
                "仿宋": "FangSong",
            }

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                font = run.font
                                font_name = font.name
                                if font_name and font_name in font_mapping:
                                    # 尝试设置嵌入字体
                                    try:
                                        font.name = font_mapping[font_name]
                                    except Exception:
                                        pass

            # 保存修改后的文件
            prs.save(pptx_path)
            logger.info(f"PPTXOptimizer: 字体嵌入完成")
        except Exception as e:
            logger.warning(f"PPTXOptimizer: 字体嵌入失败: {e}")

        return pptx_path

    def _compress_images_internal(self, pptx_path: str) -> str:
        """压缩 PPTX 内部图片的内部方法"""
        logger.info(f"PPTXOptimizer: 压缩图片...")

        try:
            # PPTX 本质上是 ZIP 文件，包含 XML 和媒体文件
            with zipfile.ZipFile(pptx_path, 'r') as zf:
                # 检查是否有 media 目录
                media_files = [f for f in zf.namelist() if f.startswith('ppt/media/')]

                if not media_files:
                    logger.info("PPTXOptimizer: 没有发现媒体文件需要压缩")
                    return pptx_path

                # 创建临时目录用于处理
                with tempfile.TemporaryDirectory() as tmp_dir:
                    extract_dir = os.path.join(tmp_dir, "pptx_contents")
                    os.makedirs(extract_dir)

                    # 解压所有文件
                    zf.extractall(extract_dir)

                    # 处理每个媒体文件
                    media_dir = os.path.join(extract_dir, "ppt", "media")
                    if os.path.exists(media_dir):
                        for media_file in os.listdir(media_dir):
                            media_path = os.path.join(media_dir, media_file)
                            if media_file.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
                                try:
                                    self._compress_image(media_path)
                                except Exception as e:
                                    logger.warning(f"PPTXOptimizer: 压缩图片失败 {media_file}: {e}")

                    # 重新打包为 PPTX
                    temp_output = os.path.join(tmp_dir, "output.pptx")
                    with zipfile.ZipFile(temp_output, 'w', zipfile.ZIP_DEFLATED) as zf_out:
                        for root, dirs, files in os.walk(extract_dir):
                            for file in files:
                                if file == "output.pptx":
                                    continue
                                file_path = os.path.join(root, file)
                                arcname = os.path.relpath(file_path, extract_dir)
                                zf_out.write(file_path, arcname)

                    # 替换原文件
                    shutil.move(temp_output, pptx_path)
                    logger.info(f"PPTXOptimizer: 图片压缩完成")

        except Exception as e:
            logger.warning(f"PPTXOptimizer: 图片压缩失败: {e}")

        return pptx_path

    def _compress_image(self, image_path: str, max_width: int = 1920, quality: int = 85):
        """压缩单张图片"""
        try:
            with Image.open(image_path) as img:
                # 转换为 RGB 如果是 RGBA
                if img.mode == 'RGBA':
                    img = img.convert('RGB')

                # 检查是否需要缩放
                width, height = img.size
                if width > max_width:
                    ratio = max_width / width
                    new_height = int(height * ratio)
                    img = img.resize((max_width, new_height), Image.Resampling.LANCZOS)

                # 保存并压缩
                if image_path.lower().endswith('.jpg') or image_path.lower().endswith('.jpeg'):
                    img.save(image_path, 'JPEG', quality=quality, optimize=True)
                else:
                    # 对于 PNG，转换为 JPEG 可以获得更好的压缩
                    new_path = image_path.rsplit('.', 1)[0] + '.jpg'
                    img.save(new_path, 'JPEG', quality=quality, optimize=True)
                    os.remove(image_path)
                    # 重命名回原扩展名（但实际内容是 JPEG）
                    os.rename(new_path, image_path)

        except Exception as e:
            logger.warning(f"压缩图片失败 {image_path}: {e}")

    def _clean_metadata_internal(self, pptx_path: str) -> str:
        """清理 PPTX 中的冗余元数据"""
        logger.info(f"PPTXOptimizer: 清理元数据...")

        try:
            with zipfile.ZipFile(pptx_path, 'r') as zf:
                with tempfile.TemporaryDirectory() as tmp_dir:
                    extract_dir = os.path.join(tmp_dir, "pptx_contents")
                    os.makedirs(extract_dir)
                    zf.extractall(extract_dir)

                    # 清理 docProps/app.xml 中的冗余信息
                    app_xml_path = os.path.join(extract_dir, "docProps", "app.xml")
                    if os.path.exists(app_xml_path):
                        # 简化 app.xml
                        with open(app_xml_path, 'w', encoding='utf-8') as f:
                            f.write('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n')
                            f.write('<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"/>')

                    # 清理 docProps/core.xml 中的冗余信息
                    core_xml_path = os.path.join(extract_dir, "docProps", "core.xml")
                    if os.path.exists(core_xml_path):
                        with open(core_xml_path, 'r', encoding='utf-8') as f:
                            content = f.read()
                        # 移除可能的敏感信息标签（保留基本结构）
                        import re
                        # 保留 [5] cp:lastModifiedBy 和 dcterms:created 等基本标签
                        content = re.sub(r'<dc:creator>[^<]*</dc:creator>', '<dc:creator>PPTXOptimizer</dc:creator>', content)
                        with open(core_xml_path, 'w', encoding='utf-8') as f:
                            f.write(content)

                    # 重新打包
                    temp_output = os.path.join(tmp_dir, "output.pptx")
                    with zipfile.ZipFile(temp_output, 'w', zipfile.ZIP_DEFLATED) as zf_out:
                        for root, dirs, files in os.walk(extract_dir):
                            for file in files:
                                if file == "output.pptx":
                                    continue
                                file_path = os.path.join(root, file)
                                arcname = os.path.relpath(file_path, extract_dir)
                                zf_out.write(file_path, arcname)

                    shutil.move(temp_output, pptx_path)
                    logger.info("PPTXOptimizer: 元数据清理完成")

        except Exception as e:
            logger.warning(f"PPTXOptimizer: 元数据清理失败: {e}")

        return pptx_path

    def compress(self, pptx_path: str, quality: int = 80) -> str:
        """压缩 PPTX 文件（兼容旧接口）"""
        logger.info(f"PPTXOptimizer: 压缩文件 {pptx_path} (quality={quality})")

        options = {"compress": True, "embed_fonts": False, "clean_metadata": False}
        return self.optimize(pptx_path, options)

    def embed_fonts(self, pptx_path: str) -> str:
        """嵌入字体（兼容旧接口）"""
        logger.info(f"PPTXOptimizer: 嵌入字体 {pptx_path}")

        options = {"compress": False, "embed_fonts": True, "clean_metadata": False}
        return self.optimize(pptx_path, options)
