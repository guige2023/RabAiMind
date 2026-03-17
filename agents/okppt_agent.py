"""
OkPPT 转换 Agent (OkPPT Agent)

负责调用 mcp-server-okppt 转换 SVG 为 PPTX
"""

import os
import json
import subprocess
from typing import Dict, Any, List, Optional
from dataclasses import dataclass
from enum import Enum
import tempfile
import shutil


class ConversionStatus(Enum):
    """转换状态"""
    PENDING = "pending"
    RUNNING = "running"
    COMPLETED = "completed"
    FAILED = "failed"


@dataclass
class ConversionResult:
    """转换结果"""
    status: ConversionStatus
    pptx_path: Optional[str] = None
    slide_count: int = 0
    error: Optional[str] = None
    metadata: Optional[Dict] = None


class OkPPTClient:
    """OkPPT MCP 客户端"""

    def __init__(self, server_path: Optional[str] = None):
        self.server_path = server_path or "/usr/local/bin/mcp-server-okppt"

    def is_available(self) -> bool:
        """检查 OkPPT 服务是否可用"""
        return os.path.exists(self.server_path)

    def convert_svg_to_pptx(
        self,
        svg_paths: List[str],
        template: Optional[str] = None,
        output_path: Optional[str] = None
    ) -> ConversionResult:
        """
        将 SVG 转换为 PPTX

        Args:
            svg_paths: SVG 文件路径列表
            template: PPT 模板路径
            output_path: 输出文件路径

        Returns:
            转换结果
        """
        if not self.is_available():
            return ConversionResult(
                status=ConversionStatus.FAILED,
                error="mcp-server-okppt 不可用"
            )

        # 创建临时目录
        temp_dir = tempfile.mkdtemp()

        try:
            # 复制 SVG 文件到临时目录
            temp_svg_paths = []
            for i, svg_path in enumerate(svg_paths):
                if os.path.exists(svg_path):
                    dest = os.path.join(temp_dir, f"slide_{i+1}.svg")
                    shutil.copy(svg_path, dest)
                    temp_svg_paths.append(dest)

            # 构建命令
            cmd = [
                self.server_path,
                "convert",
                "--input", temp_dir,
                "--output", output_path or os.path.join(temp_dir, "output.pptx")
            ]

            if template:
                cmd.extend(["--template", template])

            # 执行转换
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=300
            )

            if result.returncode == 0:
                # 找到生成的 PPTX 文件
                output_pptx = output_path or os.path.join(temp_dir, "output.pptx")

                if os.path.exists(output_pptx):
                    return ConversionResult(
                        status=ConversionStatus.COMPLETED,
                        pptx_path=output_pptx,
                        slide_count=len(temp_svg_paths),
                        metadata={"temp_dir": temp_dir}
                    )

            return ConversionResult(
                status=ConversionStatus.FAILED,
                error=result.stderr or "转换失败"
            )

        except subprocess.TimeoutExpired:
            return ConversionResult(
                status=ConversionStatus.FAILED,
                error="转换超时"
            )
        except Exception as e:
            return ConversionResult(
                status=ConversionStatus.FAILED,
                error=str(e)
            )
        finally:
            # 清理临时目录（可选，保留用于调试）
            # shutil.rmtree(temp_dir, ignore_errors=True)
            pass


class OkPPTAgent:
    """
    OkPPT 转换 Agent

    负责:
    - 调用 mcp-server-okppt 转换
    - SVG 插入 PPT 模板
    - 生成 PPTX
    """

    def __init__(self, config: Optional[Dict] = None):
        self.config = config or {}
        self.client = OkPPTClient(
            server_path=self.config.get("okppt_server_path")
        )
        self.output_dir = self.config.get("output_dir", "./output")
        self.template_dir = self.config.get("template_dir", "./templates")

    def convert(
        self,
        slides: List[Dict[str, Any]],
        template: Optional[str] = None,
        output_path: Optional[str] = None
    ) -> ConversionResult:
        """
        转换幻灯片为 PPTX

        Args:
            slides: 幻灯片列表，每项应包含 svg_path
            template: 模板路径
            output_path: 输出文件路径

        Returns:
            转换结果
        """
        # 提取 SVG 路径
        svg_paths = []
        for slide in slides:
            svg_path = slide.get("svg_path")
            if svg_path and os.path.exists(svg_path):
                svg_paths.append(svg_path)

        if not svg_paths:
            return ConversionResult(
                status=ConversionStatus.FAILED,
                error="没有有效的 SVG 文件"
            )

        # 确保输出目录存在
        os.makedirs(self.output_dir, exist_ok=True)

        # 设置输出路径
        if not output_path:
            output_path = os.path.join(self.output_dir, "presentation.pptx")

        # 优先使用 mcp-server-okppt，如果不可用则使用 python-pptx
        if self.client.is_available():
            return self.client.convert_svg_to_pptx(
                svg_paths=svg_paths,
                template=template,
                output_path=output_path
            )
        else:
            # 使用 python-pptx 后备方案
            print("mcp-server-okppt 不可用，使用 python-pptx 后备方案...")
            return self._convert_with_python_pptx(slides, output_path)

    def _convert_with_python_pptx(
        self,
        slides: List[Dict[str, Any]],
        output_path: str
    ) -> ConversionResult:
        """使用 python-pptx 转换"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation()

            # 设置幻灯片大小为 16:9
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            for slide_data in slides:
                # 添加空白幻灯片
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)

                # 添加标题
                title_text = slide_data.get("title", "")
                if title_text:
                    title_box = slide.shapes.add_textbox(
                        Inches(0.5),
                        Inches(0.3),
                        Inches(12),
                        Inches(1)
                    )
                    text_frame = title_box.text_frame
                    text_frame.text = title_text
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(44)
                        paragraph.font.bold = True

                # 添加内容
                content = slide_data.get("content", [])
                if isinstance(content, str):
                    content = content.split("\n")

                if content:
                    content_box = slide.shapes.add_textbox(
                        Inches(0.5),
                        Inches(1.5),
                        Inches(12),
                        Inches(5)
                    )
                    text_frame = content_box.text_frame
                    text_frame.word_wrap = True

                    for i, item in enumerate(content[:6]):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        p.text = f"• {item}"
                        p.font.size = Pt(24)

            prs.save(output_path)

            return ConversionResult(
                status=ConversionStatus.COMPLETED,
                pptx_path=output_path,
                slide_count=len(slides)
            )

        except Exception as e:
            return ConversionResult(
                status=ConversionStatus.FAILED,
                error=f"python-pptx 转换失败: {str(e)}"
            )

    def convert_single(
        self,
        svg_path: str,
        output_path: Optional[str] = None
    ) -> ConversionResult:
        """
        转换单个 SVG 为 PPTX（单页）

        Args:
            svg_path: SVG 文件路径
            output_path: 输出文件路径

        Returns:
            转换结果
        """
        return self.convert([{"svg_path": svg_path}], output_path=output_path)

    def add_slides(
        self,
        existing_pptx: str,
        new_slides: List[Dict[str, Any]],
        output_path: Optional[str] = None
    ) -> ConversionResult:
        """
        向现有 PPTX 添加幻灯片

        Args:
            existing_pptx: 现有 PPTX 文件路径
            new_slides: 新幻灯片列表
            output_path: 输出文件路径

        Returns:
            转换结果
        """
        # 提取新幻灯片的 SVG
        svg_paths = []
        for slide in new_slides:
            svg_path = slide.get("svg_path")
            if svg_path and os.path.exists(svg_path):
                svg_paths.append(svg_path)

        if not svg_paths:
            return ConversionResult(
                status=ConversionStatus.FAILED,
                error="没有有效的 SVG 文件"
            )

        # 由于 mcp-server-okppt 可能不支持追加，
        # 这里使用一个简化的实现：合并所有 SVG
        # 实际实现可能需要使用 python-pptx 库

        # 暂时返回错误，实际需要使用其他方法
        return ConversionResult(
            status=ConversionStatus.FAILED,
            error="暂不支持向现有 PPTX 添加幻灯片"
        )

    def get_available_templates(self) -> List[str]:
        """获取可用的模板列表"""
        templates = []

        if not os.path.exists(self.template_dir):
            return templates

        for file in os.listdir(self.template_dir):
            if file.endswith((".pptx", ".ppt")):
                templates.append(os.path.join(self.template_dir, file))

        return templates


class FallbackPPTXGenerator:
    """后备 PPTX 生成器（使用 python-pptx）"""

    def __init__(self):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            self.pptx_available = True
            self.Presentation = Presentation
            self.Inches = Inches
            self.Pt = Pt
        except ImportError:
            self.pptx_available = False

    def generate(
        self,
        slides: List[Dict[str, Any]],
        output_path: str,
        title: str = "Presentation"
    ) -> bool:
        """生成 PPTX"""
        if not self.pptx_available:
            print("警告: python-pptx 不可用，使用简化输出")
            return False

        try:
            prs = self.Presentation()

            # 设置幻灯片大小为 16:9
            prs.slide_width = self.Inches(13.333)
            prs.slide_height = self.Inches(7.5)

            for slide_data in slides:
                # 添加空白幻灯片
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)

                # 添加标题
                title_text = slide_data.get("title", "")
                if title_text:
                    title_box = slide.shapes.add_textbox(
                        self.Inches(0.5),
                        self.Inches(0.3),
                        self.Inches(12),
                        self.Inches(1)
                    )
                    text_frame = title_box.text_frame
                    text_frame.text = title_text
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = self.Pt(44)
                        paragraph.font.bold = True

                # 添加内容
                content = slide_data.get("content", [])
                if isinstance(content, str):
                    content = content.split("\n")

                if content:
                    content_box = slide.shapes.add_textbox(
                        self.Inches(0.5),
                        self.Inches(1.5),
                        self.Inches(12),
                        self.Inches(5)
                    )
                    text_frame = content_box.text_frame
                    text_frame.word_wrap = True

                    for i, item in enumerate(content[:6]):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        p.text = f"• {item}"
                        p.font.size = self.Pt(24)

            prs.save(output_path)
            return True

        except Exception as e:
            print(f"生成 PPTX 失败: {e}")
            return False


def create_okppt_agent(config: Optional[Dict] = None) -> OkPPTAgent:
    """创建 OkPPT Agent"""
    return OkPPTAgent(config)
