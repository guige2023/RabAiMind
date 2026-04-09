"""
总调度 Agent (Coordinator Agent)

负责解析用户 PPT 需求，拆分任务并分配给各 Agent，监控全流程
"""

import json
import os
import time
import uuid
import logging
import tempfile
import shutil
from datetime import datetime
from typing import Dict, Any, List, Optional
from enum import Enum

from agents.volc_okppt_tools import PromptOptimizer, get_volcano_client

logger = logging.getLogger(__name__)

# 延迟导入以避免循环依赖
PPTXOptimizerAgent = None


def _get_pptx_optimizer():
    global PPTXOptimizerAgent
    if PPTXOptimizerAgent is None:
        from agents.pptx_optimizer_agent import PPTXOptimizerAgent as cls
        PPTXOptimizerAgent = cls
    return PPTXOptimizerAgent()


class TaskStatus(Enum):
    """任务状态"""
    PENDING = "pending"
    RUNNING = "running"
    COMPLETED = "completed"
    FAILED = "failed"
    WAITING = "waiting"


class TaskType(Enum):
    """任务类型"""
    PARSE_REQUIREMENT = "parse_requirement"
    GENERATE_CONTENT = "generate_content"
    RENDER_SVG = "render_svg"
    CONVERT_PPTX = "convert_pptx"
    OPTIMIZE_PPTX = "optimize_pptx"
    VALIDATE_QUALITY = "validate_quality"


class PPTRequest:
    """PPT 生成请求"""

    def __init__(self, request_id: str, user_request: str, options: Optional[Dict] = None):
        self.request_id = request_id
        self.user_request = user_request
        self.options = options or {}
        self.scene = self.options.get("scene")
        self.slide_count = self.options.get("slide_count", 10)
        self.style = self.options.get("style", "professional")
        self.aspect_ratio = self.options.get("aspect_ratio", "16:9")
        self.created_at = datetime.now()

    def to_dict(self) -> Dict:
        return {
            "request_id": self.request_id,
            "user_request": self.user_request,
            "options": self.options,
            "created_at": self.created_at.isoformat()
        }


class Task:
    """任务单元"""

    def __init__(
        self,
        task_id: str,
        task_type: TaskType,
        description: str,
        depends_on: Optional[List[str]] = None,
        data: Optional[Dict] = None
    ):
        self.task_id = task_id
        self.task_type = task_type
        self.description = description
        self.depends_on = depends_on or []
        self.data = data or {}
        self.status = TaskStatus.PENDING
        self.result = None
        self.error = None
        self.started_at = None
        self.completed_at = None

    def to_dict(self) -> Dict:
        return {
            "task_id": self.task_id,
            "task_type": self.task_type.value,
            "description": self.description,
            "depends_on": self.depends_on,
            "status": self.status.value,
            "result": self.result,
            "error": self.error,
            "started_at": self.started_at.isoformat() if self.started_at else None,
            "completed_at": self.completed_at.isoformat() if self.completed_at else None
        }


class CoordinatorAgent:
    """
    总调度 Agent

    负责整个 PPT 生成流程的协调和调度:
    1. 解析用户 PPT 需求
    2. 拆分任务并分配给各 Agent
    3. 监控全流程执行
    4. 处理异常和重试
    """

    MAX_STEPS = 30
    DEFAULT_SLIDE_COUNT = 10
    MAX_SLIDE_COUNT = 50

    def __init__(self, config: Optional[Dict] = None):
        self.name = "CoordinatorAgent"
        self.config = config or {}
        self.max_steps = self.config.get("max_steps", self.MAX_STEPS)
        self.volcano_client = get_volcano_client()
        self._tasks: Dict[str, Task] = {}
        self._task_graph: Dict[str, List[str]] = {}
        self._current_step = 0
        self._context: Dict[str, Any] = {}

    def generate_ppt(self, user_request: str, options: Optional[Dict] = None) -> Dict[str, Any]:
        """
        生成 PPT 的主入口

        Args:
            user_request: 用户需求描述
            options: 生成选项

        Returns:
            生成结果，包含 request_id, status, slides 等
        """
        # 创建请求
        request_id = str(uuid.uuid4())
        ppt_request = PPTRequest(request_id, user_request, options)

        # 验证输入
        if not user_request or len(user_request.strip()) < 5:
            return {
                "success": False,
                "error": "需求描述太短，请提供更详细的内容",
                "request_id": request_id
            }

        # 限制幻灯片数量
        if ppt_request.slide_count > self.MAX_SLIDE_COUNT:
            ppt_request.slide_count = self.MAX_SLIDE_COUNT

        # 开始执行流程
        try:
            result = self._execute_workflow(ppt_request)
            return {
                "success": True,
                "request_id": request_id,
                "status": "completed",
                "result": result
            }
        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "request_id": request_id,
                "status": "failed"
            }

    def _execute_workflow(self, ppt_request: PPTRequest) -> Dict[str, Any]:
        """执行完整的工作流程"""
        # 重置状态
        self._tasks.clear()
        self._task_graph.clear()
        self._current_step = 0
        self._context.clear()

        # 存储请求信息
        self._context["request"] = ppt_request.to_dict()

        # 第一阶段: 解析需求
        slides = self._parse_requirement(ppt_request)

        # 第二阶段: 生成内容
        slides = self._generate_content(slides)

        # 第三阶段: 渲染 SVG
        slides = self._render_svgs(slides)

        # 第四阶段: 转换为 PPTX
        pptx_path = self._convert_to_pptx(slides)

        # 第五阶段: 优化 PPTX
        optimized_path = self._optimize_pptx(pptx_path)

        return {
            "slides": slides,
            "pptx_path": optimized_path,
            "total_slides": len(slides)
        }

    def _parse_requirement(self, ppt_request: PPTRequest) -> List[Dict]:
        """解析用户需求，生成幻灯片结构"""
        self._current_step += 1
        logger.info(f"[Step {self._current_step}] 解析需求: {ppt_request.user_request}")

        # 使用提示词优化器生成结构
        optimized_prompt = PromptOptimizer.optimize_for_ppt(
            ppt_request.user_request,
            scene=ppt_request.scene,
            slide_count=ppt_request.slide_count,
            style=ppt_request.style
        )

        # 调用火山云生成
        result = self.volcano_client.generate_text(
            prompt=optimized_prompt,
            system_prompt="你是一位专业的 PPT 内容策划专家，擅长生成结构清晰、内容丰富的 PPT 内容。",
            temperature=0.7,
            max_tokens=16000
        )

        # 解析返回的内容
        content = result["content"]
        slides = self._parse_slides_from_response(content, ppt_request.user_request, ppt_request.slide_count)
        # DEBUG removed

        self._context["slides"] = slides
        return slides

    def _parse_slides_from_response(self, content: str, user_request: str = "", slide_count: int = 7) -> List[Dict]:
        """从 API 响应中解析幻灯片结构"""

        # DEBUG removed

        # 尝试提取 JSON
        json_match = None
        in_json_block = False
        for line in content.split("\n"):
            if "```json" in line:
                in_json_block = True
                json_match = []
                continue
            elif "```" in line and in_json_block:
                in_json_block = False
                break
            if in_json_block and line.strip():
                json_match.append(line)

        if json_match:
            try:
                data = json.loads("\n".join(json_match))
                slides = data.get("slides", [])
                if slides:
                    # DEBUG removed
                    return slides[:slide_count]  # 限制数量
            except json.JSONDecodeError as e:
                pass

        # 尝试直接解析整个响应为JSON
        try:
            data = json.loads(content)
            slides = data.get("slides", [])
            if slides:
                # DEBUG removed
                return slides[:slide_count]
        except json.JSONDecodeError:
            pass

        # 如果无法解析 JSON，打印原始响应供调试
        # DEBUG removed

        # 尝试按行解析标题
        lines = [l.strip() for l in content.split("\n") if l.strip()]
        slides = []
        current_slide = None

        for line in lines:
            # 数字开头的行作为新幻灯片标题
            if line and (line[0].isdigit() or line.startswith("#") or "页" in line or "slide" in line.lower()):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {"type": "content", "title": line.lstrip("0123456789.# ").strip(), "content": []}
            elif current_slide and len(line) > 5:
                current_slide["content"].append(line)

        if current_slide:
            slides.append(current_slide)

        if slides:
            # DEBUG removed
            return slides[:slide_count]

        # 如果仍然无法解析，使用 LLM 响应内容作为基础
        # DEBUG removed
        # 提取前几行作为标题
        key_lines = [l.strip() for l in content.split("\n") if len(l.strip()) > 10][:slide_count]
        slides = []
        for i, line in enumerate(key_lines):
            slides.append({
                "type": "content",
                "title": f"第{i+1}页",
                "content": [line] if line else ["内容"]
            })

        # 如果还是不够，补充默认幻灯片
        while len(slides) < slide_count:
            slides.append({
                "type": "content",
                "title": f"第{len(slides)+1}页",
                "content": ["继续完善中..."]
            })

        # DEBUG removed
        return slides

    def _generate_content(self, slides: List[Dict]) -> List[Dict]:
        """生成每页的详细内容"""
        self._current_step += 1
        logger.info(f"[Step {self._current_step}] 生成内容，共 {len(slides)} 页")

        for i, slide in enumerate(slides):
            self._current_step += 1
            if self._current_step > self.max_steps:
                raise Exception(f"超过最大步骤数 {self.max_steps}")

            logger.info(f"[Step {self._current_step}] 生成第 {i+1} 页内容: {slide.get('title', 'untitled')}")

            # 为每一页生成更详细的内容
            prompt = f"""为 PPT 页面生成详细内容。

页面标题: {slide.get('title', '无标题')}
页面类型: {slide.get('type', 'content')}

请生成 3-5 个关键要点，用简洁的语言描述。"""

            result = self.volcano_client.generate_text(
                prompt=prompt,
                temperature=0.7,
                max_tokens=2000
            )

            # 更新幻灯片内容
            slide["generated_content"] = result["content"]

        return slides

    def _render_svgs(self, slides: List[Dict]) -> List[Dict]:
        """渲染 SVG"""
        self._current_step += 1
        logger.info(f"[Step {self._current_step}] 渲染 SVG，共 {len(slides)} 页")

        # 创建临时输出目录
        output_dir = self._context.get("output_dir") or os.path.join(tempfile.gettempdir(), "rabai_svg_render")
        os.makedirs(output_dir, exist_ok=True)
        self._context["output_dir"] = output_dir

        # 尝试导入 cairosvg 进行 SVG 渲染
        try:
            import cairosvg
            HAS_CAIRO = True
        except ImportError:
            HAS_CAIRO = False
            logger.warning("cairosvg 未安装，SVG 渲染将使用备用方案")

        for i, slide in enumerate(slides):
            self._current_step += 1
            if self._current_step > self.max_steps:
                raise Exception(f"超过最大步骤数 {self.max_steps}")

            logger.info(f"[Step {self._current_step}] 渲染第 {i+1} 页 SVG")

            # 生成 SVG 内容
            svg_content = self._generate_svg_for_slide(slide, i + 1)

            # 保存 SVG 文件
            svg_path = os.path.join(output_dir, f"slide_{i+1}.svg")
            with open(svg_path, 'w', encoding='utf-8') as f:
                f.write(svg_content)

            slide["svg_path"] = svg_path

            # 如果 cairosvg 可用，渲染 PNG
            if HAS_CAIRO:
                try:
                    png_path = os.path.join(output_dir, f"slide_{i+1}.png")
                    cairosvg.svg2png(
                        bytestring=svg_content.encode('utf-8'),
                        write_to=png_path,
                        output_width=1920,
                        output_height=1080
                    )
                    slide["png_path"] = png_path
                    logger.info(f"[Step {self._current_step}] SVG 渲染为 PNG: {png_path}")
                except Exception as e:
                    logger.warning(f"SVG 转 PNG 失败: {e}")
                    slide["png_path"] = None
            else:
                slide["png_path"] = None

        return slides

    def _generate_svg_for_slide(self, slide: Dict, page_num: int) -> str:
        """为单页幻灯片生成 SVG 内容"""
        title = slide.get("title", f"第 {page_num} 页")
        content = slide.get("generated_content", "") or slide.get("content", [])

        # 处理内容为文本
        if isinstance(content, list):
            content_text = "\n".join(str(c) for c in content)
        else:
            content_text = str(content)

        # 生成 SVG
        svg_width = 1920
        svg_height = 1080
        margin = 80
        title_y = 120
        content_y = 220

        # 转义 HTML 特殊字符
        def escape_html(text):
            return (text.replace("&", "&amp;")
                     .replace("<", "&lt;")
                     .replace(">", "&gt;")
                     .replace('"', "&quot;"))

        # 分行处理内容
        lines = content_text.split('\n')[:15]  # 最多15行

        content_svg = ""
        for j, line in enumerate(lines):
            y = content_y + j * 50
            content_svg += f'<text x="{margin}" y="{y}" font-family="Microsoft YaHei, SimHei, sans-serif" font-size="28" fill="#333">{escape_html(line)}</text>'

        svg = f'''<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" width="{svg_width}" height="{svg_height}" viewBox="0 0 {svg_width} {svg_height}">
  <defs>
    <linearGradient id="bg_grad" x1="0%" y1="0%" x2="100%" y2="100%">
      <stop offset="0%" style="stop-color:#f8f9fa;stop-opacity:1" />
      <stop offset="100%" style="stop-color:#e9ecef;stop-opacity:1" />
    </linearGradient>
  </defs>
  <rect width="100%" height="100%" fill="url(#bg_grad)"/>
  <rect x="{margin}" y="60" width="200" height="6" rx="3" fill="#4A90E2"/>
  <text x="{margin}" y="{title_y}" font-family="Microsoft YaHei, SimHei, sans-serif" font-size="48" font-weight="bold" fill="#1a1a1a">{escape_html(title)}</text>
  <line x1="{margin}" y1="160" x2="{svg_width - margin}" y2="160" stroke="#dee2e6" stroke-width="2"/>
  {content_svg}
  <text x="{svg_width - margin - 150}" y="{svg_height - 50}" font-family="Arial, sans-serif" font-size="20" fill="#adb5bd">第 {page_num} 页</text>
</svg>'''
        return svg

    def _convert_to_pptx(self, slides: List[Dict]) -> str:
        """转换为 PPTX"""
        self._current_step += 1
        logger.info(f"[Step {self._current_step}] 转换为 PPTX，共 {len(slides)} 页")

        # 创建输出目录
        output_dir = self._context.get("output_dir") or os.path.join(tempfile.gettempdir(), "rabai_pptx")
        os.makedirs(output_dir, exist_ok=True)

        # 导入 python-pptx
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError("python-pptx is required: pip install python-pptx")

        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)

        # 设置空白布局
        blank_layout = prs.slide_layouts[6]  # 通常是空白布局

        output_pptx = os.path.join(output_dir, "presentation.pptx")

        for i, slide in enumerate(slides):
            self._current_step += 1
            if self._current_step > self.max_steps:
                raise Exception(f"超过最大步骤数 {self.max_steps}")

            # 添加幻灯片
            ppt_slide = prs.slides.add_slide(blank_layout)

            title = slide.get("title", f"第 {i+1} 页")
            content = slide.get("generated_content", "") or slide.get("content", [])

            # 处理内容
            if isinstance(content, list):
                content_text = "\n".join(str(c) for c in content)
            else:
                content_text = str(content)

            # 添加标题
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(12.333)
            height = Inches(1)

            title_box = ppt_slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.text = title
            title_para.font.size = Pt(40)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(26, 26, 26)

            # 添加内容
            content_top = Inches(1.8)
            content_height = Inches(5)

            content_box = ppt_slide.shapes.add_textbox(left, content_top, width, content_height)
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            lines = content_text.split('\n')
            for j, line in enumerate(lines[:20]):  # 最多20行
                if j == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()
                para.text = line.strip() if line.strip() else " "
                para.font.size = Pt(24)
                para.font.color.rgb = RGBColor(51, 51, 51)
                para.space_after = Pt(12)

            # 如果有 PNG 路径，添加图片
            png_path = slide.get("png_path")
            if png_path and os.path.exists(png_path):
                try:
                    img_left = Inches(9)
                    img_top = Inches(1)
                    img_width = Inches(3.5)
                    img_height = Inches(2)
                    ppt_slide.shapes.add_picture(png_path, img_left, img_top, width=img_width, height=img_height)
                except Exception as e:
                    logger.warning(f"添加图片失败: {e}")

            logger.info(f"[Step {self._current_step}] 生成第 {i+1} 页 PPTX")

        # 保存
        prs.save(output_pptx)
        logger.info(f"PPTX 生成完成: {output_pptx}")

        self._context["pptx_path"] = output_pptx
        return output_pptx

    def _optimize_pptx(self, pptx_path: str) -> str:
        """优化 PPTX"""
        self._current_step += 1
        logger.info(f"[Step {self._current_step}] 优化 PPTX: {pptx_path}")

        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PPTX 文件不存在: {pptx_path}")

        # 调用 PPTXOptimizerAgent 进行真实优化
        try:
            optimizer = _get_pptx_optimizer()
            optimized_path = optimizer.optimize(
                pptx_path,
                options={
                    "embed_fonts": True,
                    "compress": True,
                    "clean_metadata": True
                }
            )
            logger.info(f"[Step {self._current_step}] PPTX 优化完成: {optimized_path}")
            return optimized_path
        except Exception as e:
            logger.warning(f"PPTX 优化失败: {e}，使用原文件")
            return pptx_path

    def get_task_status(self, request_id: str) -> Dict[str, Any]:
        """获取任务状态"""
        return {
            "request_id": request_id,
            "status": "completed",
            "progress": 100,
            "current_step": self._current_step,
            "tasks": [task.to_dict() for task in self._tasks.values()]
        }

    def cancel_task(self, request_id: str) -> bool:
        """取消任务"""
        # 取消逻辑
        return True


def create_coordinator(config: Optional[Dict] = None) -> CoordinatorAgent:
    """创建调度器实例"""
    return CoordinatorAgent(config)
